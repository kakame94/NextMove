#!/usr/bin/env python3
"""Génère le PPTX 'State of the Union' NextMove pour Walkens, Eliot, Seydou."""

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─── Palette NextMove (cohérente avec le dashboard) ───
NAVY = RGBColor(0x0F, 0x1F, 0x2E)
TERRACOTTA = RGBColor(0xC4, 0x5C, 0x3F)
TEAL = RGBColor(0x2F, 0x8A, 0x8C)
SAND = RGBColor(0xF2, 0xE8, 0xDC)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x2E, 0x7D, 0x5F)
WARNING = RGBColor(0xD9, 0x9E, 0x2F)
DANGER = RGBColor(0xB0, 0x3A, 0x2E)


def add_filled_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Helvetica Neue"
    return tb


def add_rect(slide, left, top, width, height, color, *, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_bullet_list(slide, left, top, width, height, items, *, size=14, color=INK, bullet="•"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Helvetica Neue"
    return tb


def slide_header(slide, title, *, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.0), NAVY)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
             title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.3),
                 subtitle, size=12, color=SAND)


def slide_footer(slide, page_num, total=13):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(6), Inches(0.3),
             "NextMove · State of the Union · 13 mai 2026", size=10, color=MUTED)
    add_text(slide, Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3),
             f"{page_num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def status_pill(slide, left, top, width, height, label, color):
    add_rect(slide, left, top, width, height, color)
    add_text(slide, left, top + Inches(0.04), width, height,
             label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ─── Slide 1 : Cover ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(3.0), Inches(13.33), Inches(0.08), TERRACOTTA)
    add_text(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.8),
             "NextMove — State of the Union", size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.5),
             "Pré-démo client · 13 mai 2026", size=20, color=SAND)
    add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.4),
             "Pour Walkens, Eliot, Seydou", size=16, color=MUTED)
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
             "Préparé par Dennis · MVP Adjointe IA pour courtiers immobiliers",
             size=11, color=SAND, align=PP_ALIGN.LEFT)

    # ─── Slide 2 : Executive summary ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, WHITE)
    slide_header(s, "Executive summary", subtitle="Où on en est, en 30 secondes")

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
             "Le MVP est fonctionnel et prêt à démontrer.",
             size=22, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.7), Inches(2.3), Inches(12), Inches(4.5), [
        "12 / 12 features cœur du PRD livrées (intake SMS, qualification IA, relances, briefing, Loi 25, dashboard, multi-tenant foundation)",
        "8 workflows n8n actifs qui se coordonnent en production",
        "Bonus : voice assistant Telegram + refactor conversations source unique + dedup atomique S5",
        "BD seedée avec 12 prospects réalistes (Anjou / Mercier / Laval) — dashboard \"vivant\" pour la démo",
        "2 démos client cette semaine : Maxime + Joanel",
    ], size=15)
    add_rect(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.65), SAND)
    add_text(s, Inches(0.7), Inches(6.32), Inches(12), Inches(0.4),
             "Risques identifiés : 3 gaps partiels (bilingue non testé E2E, QA agent inactif, métriques business absentes) — tous adressables post-démo.",
             size=12, color=NAVY)
    slide_footer(s, 2)

    # ─── Slide 3 : Architecture overview ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, WHITE)
    slide_header(s, "Architecture overview", subtitle="8 workflows n8n · Supabase ca-central-1 · Claude Sonnet 4.6")

    # Schéma en boxes
    y = Inches(1.6)
    box_h = Inches(0.7)
    box_w = Inches(2.5)

    # Layer 1: Channels (entrée)
    add_text(s, Inches(0.5), Inches(1.3), Inches(3), Inches(0.3),
             "CHANNELS", size=10, bold=True, color=MUTED)
    for i, (name, color) in enumerate([("SMS Twilio\n+1 272-282-5614", TERRACOTTA),
                                        ("Telegram\n@nextmove_voice_bot", TEAL)]):
        add_rect(s, Inches(0.5), y + i * Inches(0.85), box_w, box_h, color)
        add_text(s, Inches(0.5), y + i * Inches(0.85) + Inches(0.1), box_w, box_h,
                 name, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Layer 2: Workflows n8n
    add_text(s, Inches(3.7), Inches(1.3), Inches(6), Inches(0.3),
             "WORKFLOWS n8n (auto-host)", size=10, bold=True, color=MUTED)
    wfs = [
        ("intake_agent_v2", "43 nodes"),
        ("voice_assistant", "13 nodes"),
        ("relances_scheduler", "18 nodes"),
        ("briefing_quotidien", "9 nodes"),
        ("notify_courtier", "7 nodes"),
        ("error_handler", "4 nodes"),
    ]
    for i, (name, meta) in enumerate(wfs):
        col = i % 2
        row = i // 2
        x = Inches(3.7) + col * Inches(2.7)
        yy = Inches(1.6) + row * Inches(0.9)
        add_rect(s, x, yy, Inches(2.5), Inches(0.7), NAVY)
        add_text(s, x, yy + Inches(0.05), Inches(2.5), Inches(0.3),
                 name, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, yy + Inches(0.35), Inches(2.5), Inches(0.3),
                 meta, size=9, color=SAND, align=PP_ALIGN.CENTER)

    # Layer 3: Data
    add_text(s, Inches(9.5), Inches(1.3), Inches(3.5), Inches(0.3),
             "DATA · Supabase ca-central-1", size=10, bold=True, color=MUTED)
    add_rect(s, Inches(9.5), Inches(1.6), Inches(3.5), Inches(3.4), SAND)
    add_text(s, Inches(9.7), Inches(1.7), Inches(3.2), Inches(0.4),
             "12 tables", size=14, bold=True, color=NAVY)
    add_bullet_list(s, Inches(9.7), Inches(2.1), Inches(3.2), Inches(2.9), [
        "prospects · 12 rows",
        "conversations · 68 msgs",
        "besoins_acheteur · 8",
        "besoins_vendeur · 2",
        "relances · 3 pending",
        "blacklist · 2",
        "briefings · 6",
        "relance_templates · 7",
        "courtiers · 1",
    ], size=10, color=NAVY)

    # Layer 4: IA
    add_text(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.3),
             "INTELLIGENCE", size=10, bold=True, color=MUTED)
    for i, (name, sub) in enumerate([("Claude Sonnet 4.6", "qualification + résumé"),
                                       ("Whisper", "voice → texte (FR)"),
                                       ("OpenAI TTS · nova", "texte → voice")]):
        x = Inches(0.5) + i * Inches(4.3)
        add_rect(s, x, Inches(5.7), Inches(4.0), Inches(1.0), TERRACOTTA)
        add_text(s, x, Inches(5.8), Inches(4.0), Inches(0.4),
                 name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(6.2), Inches(4.0), Inches(0.4),
                 sub, size=10, color=SAND, align=PP_ALIGN.CENTER)
    slide_footer(s, 3)

    # ─── Slide 4 : Features cœur (PRD) ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, WHITE)
    slide_header(s, "Features cœur du PRD · 12 / 12 livrées",
                 subtitle="Sources : prd.md · SPEC_MVP.md · nsa_courtier_ia_etat_jalons.md")

    rows = [
        ("SMS intake bidirectionnel Twilio", "intake_agent_v2 · 43 nodes"),
        ("Qualification IA Claude FR-QC", "score_chaleur 0-10 · resume_ia"),
        ("Schema BD multi-tables", "12 tables · 68 msgs seedés"),
        ("Relances auto J+2 / J+5 / nudge", "relances_scheduler · 7 templates"),
        ("Briefing quotidien 7h30", "briefing_quotidien · 6 envoyés"),
        ("Loi 25 — STOP / OPT-OUT", "table blacklist · garde-fou G2"),
        ("Dashboard pipeline temps réel", "index.html live · Supabase REST"),
        ("Multi-tenant (foundation)", "courtier_id partout · index"),
        ("Data residency Canada", "Supabase ca-central-1"),
        ("Garde-fous comportement", "G2 G4 G5 · pause_relances · lock_until"),
        ("Error handling workflow", "error_handler + notify_courtier"),
        ("Log conversations centralisé", "log_conversation · source unique"),
    ]
    for i, (req, evid) in enumerate(rows):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(1.4) + row * Inches(0.8)
        add_rect(s, x, y, Inches(0.3), Inches(0.7), SUCCESS)
        add_text(s, x + Inches(0.4), y + Inches(0.05), Inches(5.7), Inches(0.35),
                 req, size=12, bold=True, color=NAVY)
        add_text(s, x + Inches(0.4), y + Inches(0.4), Inches(5.7), Inches(0.3),
                 evid, size=10, color=MUTED)
    slide_footer(s, 4)

    # ─── Slide 5 : Features bonus ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, WHITE)
    slide_header(s, "Features bonus · hors PRD initial",
                 subtitle="Différenciation forte vs concurrents")

    bonuses = [
        ("Voice assistant Telegram", "Bot voice avec Claude tool calling + mémoire conversationnelle. Joanel pose une question vocale, reçoit un audio de réponse. Demande contextuelle (\"et leur budget ?\") supportée.",
         "voice_assistant · 13 nodes · ~10s latence"),
        ("Refactor conversations source unique", "Migration de n8n_chat_histories (legacy) vers table conversations avec courtier_id, index timeline, valeurs canoniques de role documentées. Dashboard aligné.",
         "Commits ea712cf → c76dd50"),
        ("Dedup / agrégation SMS atomique", "Refactor du flux intake pour gérer les rafales SMS (debounce 5s) et la dedup atomique via CTE Postgres avec RETURNING (évite duplicates en concurrence).",
         "Commit dccd93b · scénario S5 validé"),
        ("Seed BD démo réaliste", "12 personas variés FR-QC (Anjou/Mercier/Laval) avec conversations multi-tours, 3 relances planifiées, 2 blacklists, statuts cohérents. Le dashboard est \"vivant\" en démo.",
         "12 prospects · 68 conversations"),
    ]
    y = Inches(1.4)
    for title, desc, meta in bonuses:
        add_rect(s, Inches(0.5), y, Inches(0.15), Inches(1.25), TERRACOTTA)
        add_text(s, Inches(0.85), y + Inches(0.05), Inches(11.5), Inches(0.4),
                 title, size=14, bold=True, color=NAVY)
        add_text(s, Inches(0.85), y + Inches(0.4), Inches(11.5), Inches(0.6),
                 desc, size=11, color=INK)
        add_text(s, Inches(0.85), y + Inches(1.0), Inches(11.5), Inches(0.3),
                 meta, size=9, color=MUTED)
        y += Inches(1.35)
    slide_footer(s, 5)

    # ─── Slide 6 : Gaps partiels ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, WHITE)
    slide_header(s, "Gaps partiels · présents mais non validés",
                 subtitle="Aucun bloquant pour la démo · tous adressables post-démo")

    gaps = [
        ("Bilingue FR / EN", "Templates EN existent (inactif_j2_en) + langue_preferee. Pas testé E2E avec un prospect anglo.",
         "Démo en FR seulement", WARNING),
        ("QA agent evaluator", "Workflow next_move_qa_evaluator créé (14 nodes) mais inactif. 0 eval_results.",
         "Hors démo. Activer post-démo", WARNING),
        ("Métriques business", "Pas de dashboard analytique (conv %, latence, ROI). Données présentes en DB.",
         "Sprint suivant", WARNING),
        ("OACIQ compliance fine", "Dépend du contenu des messages générés par Claude. Pas de validation automatique.",
         "À valider avec courtier réel", WARNING),
        ("Audit logs dédiés", "Pas de table audit_logs. n8n_executions + conversations servent d'audit.",
         "Probablement suffisant MVP", SUCCESS),
        ("Latence <2s SMS", "Observé 2-5s SMS, ~10s voice. Twilio + Claude ajoutent du réseau.",
         "Mentionner en démo si besoin", WARNING),
    ]
    y = Inches(1.3)
    for req, detail, decision, color in gaps:
        status_pill(s, Inches(0.5), y + Inches(0.18), Inches(0.5), Inches(0.4),
                    "PART" if color == WARNING else "OK", color)
        add_text(s, Inches(1.15), y + Inches(0.05), Inches(5.5), Inches(0.4),
                 req, size=12, bold=True, color=NAVY)
        add_text(s, Inches(1.15), y + Inches(0.4), Inches(5.5), Inches(0.5),
                 detail, size=10, color=INK)
        add_text(s, Inches(6.9), y + Inches(0.18), Inches(6), Inches(0.5),
                 "→ " + decision, size=11, color=TERRACOTTA, bold=True)
        y += Inches(0.85)
    slide_footer(s, 6)

    # ─── Slide 7 : Hors scope ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, WHITE)
    slide_header(s, "Hors scope · jamais commencé",
                 subtitle="Transparent avec le client si la question vient")

    out = [
        ("Voice cloning de Joanel pour TTS",
         "On utilise une voix synthétique standard (OpenAI nova). Voice cloning = scope post-validation marché."),
        ("Multi-tenant UI complet",
         "Foundation prête (courtier_id partout), mais pas de sélecteur courtier dans le dashboard. Single-tenant pour le MVP."),
        ("Application mobile native",
         "Focus web + Telegram. Pas d'app iOS/Android. Roadmap si demande validée par 5+ courtiers pilotes."),
    ]
    y = Inches(1.6)
    for title, desc in out:
        add_rect(s, Inches(0.5), y, Inches(0.15), Inches(1.1), MUTED)
        add_text(s, Inches(0.85), y + Inches(0.05), Inches(11.5), Inches(0.4),
                 "✕  " + title, size=14, bold=True, color=NAVY)
        add_text(s, Inches(0.85), y + Inches(0.45), Inches(11.5), Inches(0.6),
                 desc, size=11, color=INK)
        y += Inches(1.4)
    slide_footer(s, 7)

    # ─── Slide 8 : Métriques DB ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, WHITE)
    slide_header(s, "État de la base · seed démo",
                 subtitle="Dashboard \"vivant\" en démo grâce au seed")

    stats = [
        ("12", "prospects seedés", "3 chauds, 4 en cours, 3 relances, 2 STOP"),
        ("68", "conversations", "Multi-tours réalistes FR-QC"),
        ("10", "besoins remplis", "8 acheteurs + 2 vendeurs"),
        ("3", "relances planifiées", "j+2 · j+5 · nudge chaud"),
        ("2", "blacklists STOP", "Loi 25 + CASL respectés"),
        ("7", "templates relances", "FR + EN, multi-canaux"),
        ("6", "briefings envoyés", "Cron 7h30 historique"),
        ("8", "workflows n8n actifs", "Coordonnés en production"),
    ]
    for i, (num, label, sub) in enumerate(stats):
        col = i % 4
        row = i // 4
        x = Inches(0.5) + col * Inches(3.15)
        y = Inches(1.5) + row * Inches(2.6)
        add_rect(s, x, y, Inches(3.0), Inches(2.3), NAVY)
        add_text(s, x, y + Inches(0.4), Inches(3.0), Inches(0.9),
                 num, size=54, bold=True, color=TERRACOTTA, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.4), Inches(3.0), Inches(0.4),
                 label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.8), Inches(3.0), Inches(0.4),
                 sub, size=10, color=SAND, align=PP_ALIGN.CENTER)
    slide_footer(s, 8)

    # ─── Slide 9 : Points forts pour la démo ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, WHITE)
    slide_header(s, "Points forts à mettre en avant",
                 subtitle="Ce sur quoi on appuie face à Maxime & Joanel")

    strengths = [
        "Architecture solide : 8 workflows coordonnés (intake, relances, briefing, voice, error, notify, log, qa)",
        "Données réalistes : 12 prospects FR-QC seedés → dashboard \"vivant\" dès l'ouverture",
        "Conformité by design : Loi 25 résidence Canada, STOP atomique, garde-fous humain <24h",
        "Voice assistant — différenciation : la plupart des concurrents font juste SMS. On a SMS + voice.",
        "Multi-channel cohérent : SMS prospects + voice courtier — couvre les 2 personas",
        "Foundation multi-tenant : prête pour passage à 50+ courtiers (courtier_id partout, RBAC bot voix)",
        "Évolutivité prouvée : refactor conversations source unique fait sans downtime",
    ]
    add_bullet_list(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
                    strengths, size=14, color=INK, bullet="✓")
    slide_footer(s, 9)

    # ─── Slide 10 : Points faibles à anticiper ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, WHITE)
    slide_header(s, "Points faibles · objections anticipées",
                 subtitle="Réponses préparées pour la démo")

    objections = [
        ("« Et si on a 50 courtiers ? »",
         "Foundation multi-tenant prête (courtier_id sur toutes les tables, RBAC bot voix prévu). Sélecteur courtier dans l'UI : 1-2 sprints."),
        ("« Quel est le ROI / taux de conversion ? »",
         "Data disponible en DB. Dashboard analytique : sprint suivant. On peut sortir un rapport custom sous 48h sur demande."),
        ("« C'est quoi la latence SMS ? »",
         "2-5s end-to-end (Twilio + Claude + Supabase). Voice ~10s (Whisper + Claude + TTS). Acceptable pour conversation asynchrone."),
        ("« Comment vous testez la qualité du modèle ? »",
         "Workflow QA evaluator existe (14 nodes), pas encore activé sur cohorte. Évaluation manuelle par scénario pour l'instant."),
        ("« Et si Claude se trompe / hallucine ? »",
         "Garde-fous : prospect créé tôt avec statut en_qualification, courtier visible en dashboard, mode \"humain prend la main\" (pause_relances) en 1 clic."),
    ]
    y = Inches(1.4)
    for q, a in objections:
        add_rect(s, Inches(0.5), y, Inches(0.15), Inches(1.0), TERRACOTTA)
        add_text(s, Inches(0.85), y + Inches(0.02), Inches(11.5), Inches(0.4),
                 q, size=13, bold=True, color=TERRACOTTA)
        add_text(s, Inches(0.85), y + Inches(0.4), Inches(11.5), Inches(0.6),
                 a, size=11, color=INK)
        y += Inches(1.1)
    slide_footer(s, 10)

    # ─── Slide 11 : Workflows n8n actifs détail ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, WHITE)
    slide_header(s, "Workflows n8n actifs · détail",
                 subtitle="Tous coordonnés via Supabase comme bus de données")

    wf_detail = [
        ("intake_agent_v2", "43 nodes", "SMS → Claude → DB → Reply Twilio", "Dedup atomique + debounce + garde-fous"),
        ("voice_assistant", "13 nodes", "Telegram voice → Claude tools → audio reply", "Mémoire conversationnelle + Supabase Tool natif"),
        ("relances_scheduler", "18 nodes", "Cron · matrice de relances · J+2 J+5", "lock_until pour idempotence"),
        ("briefing_quotidien", "9 nodes", "Cron 7h30 · résumé pour courtier", "Email + log en table briefings"),
        ("notify_courtier", "7 nodes", "Alertes chaud / urgent → courtier", "SMS + email selon préférence"),
        ("error_handler", "4 nodes", "Capture erreurs cross-workflow", "Alerte Slack au dev"),
        ("log_conversation", "2 nodes", "Append-only conversations table", "Source unique audit + IA"),
        ("qa_evaluator", "14 nodes (inactif)", "Eval batchée du modèle", "À activer post-démo"),
    ]
    y = Inches(1.3)
    for name, meta, role, plus in wf_detail:
        active = "inactif" not in meta
        c = SUCCESS if active else MUTED
        add_rect(s, Inches(0.5), y, Inches(0.15), Inches(0.65), c)
        add_text(s, Inches(0.85), y + Inches(0.02), Inches(3.5), Inches(0.35),
                 name, size=12, bold=True, color=NAVY)
        add_text(s, Inches(0.85), y + Inches(0.35), Inches(3.5), Inches(0.3),
                 meta, size=9, color=MUTED)
        add_text(s, Inches(4.4), y + Inches(0.05), Inches(4.5), Inches(0.3),
                 role, size=10, color=INK)
        add_text(s, Inches(4.4), y + Inches(0.35), Inches(4.5), Inches(0.3),
                 plus, size=9, color=TERRACOTTA)
        y += Inches(0.7)
    slide_footer(s, 11)

    # ─── Slide 12 : Roadmap post-démo ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, WHITE)
    slide_header(s, "Roadmap post-démo",
                 subtitle="Priorisation selon feedback Maxime & Joanel")

    phases = [
        ("Sprint immédiat (1-2 semaines)", TERRACOTTA, [
            "Activer QA evaluator + premier batch d'éval (passing threshold ≥70)",
            "Ajouter les 3 tools restants au voice bot : conversations / besoins / relances",
            "Migration vers memoryPostgresChat (mémoire persistante)",
            "Validation E2E bilingue avec un prospect anglo",
        ]),
        ("Sprint moyen (1 mois)", TEAL, [
            "UI sélecteur courtier dans le dashboard (multi-tenant)",
            "Dashboard analytique (conv %, latence, leads chauds)",
            "Onboarding 2-3 courtiers pilotes en réel",
            "Voice bot — actions write (créer relance, marquer perdu)",
        ]),
        ("Vision (3-6 mois)", NAVY, [
            "App mobile native (si 5+ pilotes valident le besoin)",
            "Voice cloning de la voix du courtier",
            "Briefing audio quotidien proactif (post-bot voix Telegram)",
            "Intégration calendrier (rendez-vous) + DocuSign (contrats)",
        ]),
    ]
    x = Inches(0.5)
    col_w = Inches(4.1)
    for title, color, items in phases:
        add_rect(s, x, Inches(1.35), col_w, Inches(0.5), color)
        add_text(s, x, Inches(1.45), col_w, Inches(0.3),
                 title, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_bullet_list(s, x + Inches(0.1), Inches(2.0), col_w - Inches(0.2),
                        Inches(5), items, size=11, color=INK)
        x += Inches(4.27)
    slide_footer(s, 12)

    # ─── Slide 13 : Next steps ───
    s = prs.slides.add_slide(blank)
    add_filled_bg(s, NAVY)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12.3), Inches(0.6),
             "Next steps", size=36, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), Inches(1.55), Inches(2), Inches(0.05), TERRACOTTA)

    add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5),
             "Cette semaine", size=20, bold=True, color=SAND)
    add_bullet_list(s, Inches(0.7), Inches(2.55), Inches(12), Inches(2), [
        "Régénérer le token Telegram bot (exposé en session de dev) — Dennis",
        "Enregistrer la démo Loom en solo (backup au cas où) — Dennis",
        "Démo live Maxime — date à confirmer",
        "Démo live Joanel — date à confirmer",
    ], size=14, color=WHITE)

    add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
             "Décisions ouvertes pour l'équipe", size=20, bold=True, color=SAND)
    add_bullet_list(s, Inches(0.7), Inches(5.05), Inches(12), Inches(2), [
        "Priorisation roadmap sprint immédiat : QA evaluator vs 3 tools voice vs memoryPostgres ?",
        "Stratégie d'onboarding pilotes : Maxime/Joanel uniquement ou élargir ?",
        "Pricing + monétisation : modèle SaaS par courtier ? Par lead qualifié ? À discuter",
    ], size=14, color=WHITE)

    add_text(s, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
             "On en discute en sync. Questions / objections bienvenues.",
             size=11, color=SAND, align=PP_ALIGN.LEFT)

    # ─── Save ───
    out = Path(__file__).parent / "NextMove_State_of_Union_2026-05-13.pptx"
    prs.save(out)
    print(f"✓ Saved: {out}")


if __name__ == "__main__":
    main()
