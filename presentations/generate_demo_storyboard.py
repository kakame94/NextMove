#!/usr/bin/env python3
"""Génère le PPTX storyboard démo NextMove (Maxime & Joanel)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0F, 0x1F, 0x2E)
TERRACOTTA = RGBColor(0xC4, 0x5C, 0x3F)
TEAL = RGBColor(0x2F, 0x8A, 0x8C)
SAND = RGBColor(0xF2, 0xE8, 0xDC)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x2E, 0x7D, 0x5F)
WARNING = RGBColor(0xD9, 0x9E, 0x2F)
DANGER = RGBColor(0xB0, 0x3A, 0x2E)


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Helvetica Neue"
    return tb


def add_rect(slide, left, top, width, height, color, *, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_bullets(slide, left, top, width, height, items, *, size=14,
                color=INK, bullet="•"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Helvetica Neue"
    return tb


def slide_header(slide, title, *, subtitle=None, accent=TERRACOTTA):
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.0), NAVY)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
             title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.3),
                 subtitle, size=12, color=SAND)
    add_rect(slide, Inches(0), Inches(1.0), Inches(0.5), Inches(0.05), accent)


def slide_footer(slide, page, total):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(6), Inches(0.3),
             "Storyboard démo NextMove · 13 mai 2026",
             size=10, color=MUTED)
    add_text(slide, Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3),
             f"{page} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def timing_banner(slide, left, top, width, time_label, *, color=TERRACOTTA):
    add_rect(slide, left, top, width, Inches(0.5), color)
    add_text(slide, left, top + Inches(0.08), width, Inches(0.4),
             time_label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def section_slide(prs, page, total, time_range, title, action_title, actions,
                  punchline, *, backup=None, color=TERRACOTTA):
    """Slide pour une section de la démo."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    slide_header(s, title, subtitle=f"⏱ {time_range}")

    # Action principale
    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
             action_title, size=18, bold=True, color=NAVY)

    # Actions / étapes
    add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.3),
             "Actions", size=11, bold=True, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(2.3), Inches(12), Inches(2.8),
                actions, size=13, color=INK)

    # Punchline
    y_punch = Inches(5.2)
    add_rect(s, Inches(0.5), y_punch, Inches(12.3), Inches(0.9), SAND)
    add_text(s, Inches(0.7), y_punch + Inches(0.08), Inches(12), Inches(0.3),
             "Phrase clé", size=10, bold=True, color=color)
    add_text(s, Inches(0.7), y_punch + Inches(0.4), Inches(12), Inches(0.5),
             f"“{punchline}”", size=13, color=NAVY, italic=True)

    # Backup
    if backup:
        y_bk = Inches(6.3)
        add_rect(s, Inches(0.5), y_bk, Inches(0.15), Inches(0.7), WARNING)
        add_text(s, Inches(0.85), y_bk + Inches(0.05), Inches(2), Inches(0.3),
                 "Si ça foire", size=10, bold=True, color=WARNING)
        add_text(s, Inches(0.85), y_bk + Inches(0.35), Inches(11.5),
                 Inches(0.35), backup, size=11, color=INK)

    slide_footer(s, page, total)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 13

    # ─── Slide 1 : Cover ───
    s = prs.slides.add_slide(blank)
    bg(s, NAVY)
    add_rect(s, Inches(0), Inches(3.0), Inches(13.33), Inches(0.08), TERRACOTTA)
    add_text(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.6),
             "Démo client NextMove", size=46, bold=True, color=WHITE)
    add_text(s, Inches(0.5), Inches(2.75), Inches(12.3), Inches(0.4),
             "Storyboard 10-12 min · live + Loom solo backup",
             size=18, color=SAND)
    add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.4),
             "Cibles : Maxime Belma · Joanel",
             size=16, color=MUTED)
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
             "Préparé par Dennis · 13 mai 2026 · MVP Adjointe IA",
             size=11, color=SAND)

    # ─── Slide 2 : Pre-flight checklist ───
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    slide_header(s, "Pre-flight · 15 min avant la démo",
                 subtitle="Vérifie matériel + stack + URL")
    # Matériel
    add_text(s, Inches(0.5), Inches(1.4), Inches(6), Inches(0.4),
             "Matériel", size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(0.7), Inches(1.85), Inches(5.8), Inches(2.5), [
        "Laptop chargé, mode « Ne pas déranger » ON",
        "Téléphone perso + Telegram sur @nextmove_voice_bot",
        "Téléphone testeuse pour SMS vers +1 XXX XXX XXXX",
        "Wi-Fi solide + plan B 4G",
    ], size=12)

    # Stack
    add_text(s, Inches(6.8), Inches(1.4), Inches(6), Inches(0.4),
             "Stack", size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(2.5), [
        "Server local lancé : python3 -m http.server 8765",
        "Workflow next_move_intake_agent_v2 → actif",
        "Workflow next_move_voice_assistant → actif",
        "Test rapide : voice msg → reply en <15s",
    ], size=12)

    # URLs
    add_rect(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3), SAND)
    add_text(s, Inches(0.7), Inches(4.65), Inches(12), Inches(0.4),
             "URLs dashboard selon la démo", size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.7), Inches(5.15), Inches(0.15), Inches(0.7), TERRACOTTA)
    add_text(s, Inches(1.0), Inches(5.15), Inches(11), Inches(0.3),
             "Démo Maxime", size=13, bold=True, color=NAVY)
    add_text(s, Inches(1.0), Inches(5.45), Inches(11), Inches(0.3),
             "http://localhost:8765/index.html?courtier=Maxime%20Belma",
             size=11, color=INK)

    add_rect(s, Inches(0.7), Inches(5.95), Inches(0.15), Inches(0.7), TEAL)
    add_text(s, Inches(1.0), Inches(5.95), Inches(11), Inches(0.3),
             "Démo Joanel (default)", size=13, bold=True, color=NAVY)
    add_text(s, Inches(1.0), Inches(6.25), Inches(11), Inches(0.3),
             "http://localhost:8765/index.html",
             size=11, color=INK)
    slide_footer(s, 2, total)

    # ─── Slide 3 : Timeline d'ensemble ───
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    slide_header(s, "Timeline d'ensemble", subtitle="6 sections sur 10-12 min")
    sections = [
        ("0:00\n→ 1:00", "Intro\nproblème", TERRACOTTA),
        ("1:00\n→ 2:30", "Tour\ndashboard", TEAL),
        ("2:30\n→ 5:00", "SMS\nLIVE", TERRACOTTA),
        ("5:00\n→ 7:30", "Voice\nLIVE", TEAL),
        ("7:30\n→ 9:00", "Robustesse\nconformité", TERRACOTTA),
        ("9:00\n→ 12:00", "Multi-canal\n& closing", TEAL),
    ]
    box_w = Inches(2.0)
    gap = Inches(0.1)
    total_w = box_w * len(sections) + gap * (len(sections) - 1)
    start_x = (Inches(13.33) - total_w) / 2
    for i, (time_lbl, label, color) in enumerate(sections):
        x = start_x + i * (box_w + gap)
        # Time banner
        add_rect(s, x, Inches(2.0), box_w, Inches(0.9), color)
        tb = s.shapes.add_textbox(x, Inches(2.05), box_w, Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, line in enumerate(time_lbl.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = WHITE
            run.font.name = "Helvetica Neue"
        # Label
        add_rect(s, x, Inches(2.9), box_w, Inches(1.6), NAVY)
        tb = s.shapes.add_textbox(x, Inches(3.2), box_w, Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, line in enumerate(label.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.size = Pt(15)
            run.font.bold = True
            run.font.color.rgb = WHITE
            run.font.name = "Helvetica Neue"

    add_text(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.4),
             "Variante Loom solo : mêmes sections, condensées en 8-10 min, pas de Q&A",
             size=12, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    slide_footer(s, 3, total)

    # ─── Slides 4-9 : Sections détaillées ───
    section_slide(prs, 4, total, "0:00 → 1:00",
                  "Intro & problème",
                  "Pose le contexte en 3 phrases. Pas de jargon.",
                  ["« Merci de prendre 10 minutes. Je te montre NextMove. »",
                   "70 % des prospects perdus si pas de réponse en 6h (étude OACIQ)",
                   "Tu ne peux pas être au téléphone 24/7. Embauche d'une adjointe = 50k$/an + risque humain",
                   "NextMove fait le job d'une adjointe humaine sur les premières 24h"],
                  "Concrètement, ça donne quoi. Regarde.")

    section_slide(prs, 5, total, "1:00 → 2:30",
                  "Tour du dashboard",
                  "Montre l'écran d'accueil — KPIs + prospects chauds",
                  ["Pointer 4 KPIs : 10 actifs, 5 chauds, 3.8 M$ pipeline, 3 relances",
                   "Cliquer sur Mathieu Simard (score 9) — détail prospect",
                   "Montrer resume_ia + besoins extraits",
                   "Cliquer « Voir la conversation » → messages naturels FR-QC"],
                  "Tu n'as rien fait sur ce prospect. L'IA l'a qualifié toute seule.",
                  color=TEAL)

    section_slide(prs, 6, total, "2:30 → 5:00",
                  "Scénario LIVE 1 — SMS",
                  "SMS testeuse → bot répond → nouveau prospect dans dashboard",
                  ["SMS depuis testeuse vers +1 XXX XXX XXXX : « Bonjour, je cherche un condo à Mercier-Est, budget 380k, on est 2 »",
                   "Pendant les 10-30s d'attente : commenter le flux Twilio → n8n → Claude → Supabase",
                   "Bot répond par SMS naturel → montrer sur le téléphone",
                   "F5 dashboard → nouveau prospect avec budget / secteur / type déjà extraits"],
                  "30 secondes après le SMS, tu as une fiche complète. Sans toi.",
                  backup="Si Twilio lag > 30s : « Le SMS passe par Twilio, ça peut prendre 30s, continuons. »")

    section_slide(prs, 7, total, "5:00 → 7:30",
                  "Scénario LIVE 2 — Voice Telegram",
                  "Interroger la base à la voix — démo mémoire conversationnelle",
                  ["Voice 1 → « Combien de prospects chauds j'ai actuellement ? »",
                   "Voice 2 follow-up → « Et lesquels cherchent à Anjou ? » (mémoire active)",
                   "Voice 3 → « C'est quoi le budget de Sophie Lavoie ? »",
                   "Téléphone à côté du laptop pour qu'on entende les réponses"],
                  "Tu interroges ta base à la voix. Pas besoin de sortir le laptop.",
                  color=TEAL,
                  backup="Si réponse foire : « Le modèle peut hésiter — on a un système d'éval continue. »")

    section_slide(prs, 8, total, "7:30 → 9:00",
                  "Robustesse & conformité",
                  "STOP, relances auto, garde-fou humain",
                  ["Marc Dubois (perdu, blacklisté) → « STOP respecté, Loi 25 + CASL obligatoires »",
                   "Patrick Côté (silence J+2, relance planifiée demain) → relance auto personnalisée",
                   "Garde-fou G4 : si toi tu réponds manuellement, l'IA se met en pause 24h",
                   "Tu gardes la main en 1 clic (pause_relances)"],
                  "Conformité by design. Toi tu gardes le contrôle.")

    section_slide(prs, 9, total, "9:00 → 12:00",
                  "Multi-canal, scale & closing",
                  "Vision multi-tenant + récap + offer",
                  ["SMS prospects + voice courtier = même base unifiée",
                   "Multi-tenant câblé en base — 50 courtiers = activation côté UI",
                   "Données au Canada (Supabase ca-central-1), Loi 25 + OACIQ OK",
                   "Récap 3 points : qualif sans toi · voice à la demande · conformité by design"],
                  "Essai 30 jours avec ton propre numéro Twilio cette semaine. Tu me dis ?",
                  color=TEAL,
                  backup="Question pricing : « Je note, je te reviens cet après-midi avec une proposition. »")

    # ─── Slide 10 : Variante Loom solo ───
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    slide_header(s, "Variante Loom solo · 8-10 min",
                 subtitle="Backup pré-enregistré + à envoyer avant la démo live")
    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
             "Mêmes sections, condensées (pas de Q&A)",
             size=14, bold=True, color=NAVY)
    rows = [
        ("0:00 → 0:45", "Intro problème (raccourci)"),
        ("0:45 → 2:00", "Tour dashboard (réduit)"),
        ("2:00 → 4:30", "SMS LIVE — 1 scénario complet"),
        ("4:30 → 6:30", "Voice Telegram — 3 questions"),
        ("6:30 → 8:00", "Robustesse (1 min par cas)"),
        ("8:00 → 8:30", "Recap court + lien vers démo live"),
    ]
    y = Inches(2.0)
    for time_lbl, desc in rows:
        add_rect(s, Inches(0.5), y, Inches(2.5), Inches(0.5), NAVY)
        add_text(s, Inches(0.5), y + Inches(0.1), Inches(2.5), Inches(0.3),
                 time_lbl, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(3.2), y + Inches(0.1), Inches(9.5), Inches(0.4),
                 desc, size=13, color=INK)
        y += Inches(0.6)

    # Tips Loom
    add_rect(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.1), SAND)
    add_text(s, Inches(0.7), Inches(5.9), Inches(12), Inches(0.3),
             "Tips Loom", size=12, bold=True, color=TERRACOTTA)
    add_bullets(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.7), [
        "Webcam ON dans le coin · Chapters pour découper · Envoyer le lien AVANT la démo live",
    ], size=11, color=INK)
    slide_footer(s, 10, total)

    # ─── Slide 11 : Backups par scénario ───
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    slide_header(s, "Backups par scénario",
                 subtitle="Si ça foire, voilà comment rattraper")
    bks = [
        ("Twilio SMS lent (>30s)", "Montrer une conversation seedée existante (Marie-Claude Tremblay)"),
        ("Voice bot ne répond pas", "Ouvrir le Loom solo dans un onglet, jouer la séquence voice"),
        ("Dashboard ne charge pas", "Onglet Supabase dashboard ouvert : montrer la DB direct"),
        ("Crash navigateur", "URL ?courtier=... pinnée dans les favoris, F5"),
        ("Question pricing imprévue", "« Bonne question, je te reviens cet après-midi avec une proposition »"),
        ("Question technique pointue", "« On a un dashboard tech pour le CTO, je te montre après si tu veux »"),
    ]
    y = Inches(1.4)
    for issue, fix in bks:
        add_rect(s, Inches(0.5), y, Inches(0.15), Inches(0.7), DANGER)
        add_text(s, Inches(0.85), y + Inches(0.05), Inches(5.5), Inches(0.35),
                 issue, size=13, bold=True, color=NAVY)
        add_text(s, Inches(6.5), y + Inches(0.05), Inches(6.3), Inches(0.65),
                 "→ " + fix, size=11, color=INK)
        y += Inches(0.82)
    slide_footer(s, 11, total)

    # ─── Slide 12 : Numéros & accès utiles ───
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    slide_header(s, "Numéros & accès utiles",
                 subtitle="À avoir sous la main pendant la démo")
    refs = [
        ("Bot Telegram", "@nextmove_voice_bot · allowlist user_id 7326149684 (toi)"),
        ("Numéro Twilio NextMove", "+1 XXX XXX XXXX (envoyer SMS ici pour démo)"),
        ("Testeuse SMS (femme)", "+1 579 421 6910"),
        ("Téléphone Dennis", "+1 XXX XXX XXXX (reçoit briefing SMS 7h30)"),
        ("Supabase project", "fhqybnkxqfvbsjvwrcob · ca-central-1"),
        ("Workflow voice (n8n)", "id 0UNAKGRciurrtlez"),
        ("Workflow intake (n8n)", "id nmmmJu6HRwq0nqyI"),
        ("Workflow briefing (n8n)", "id 9MFuSmxVnsJUQove"),
    ]
    y = Inches(1.4)
    for label, val in refs:
        add_text(s, Inches(0.5), y, Inches(4.0), Inches(0.4),
                 label, size=12, bold=True, color=TERRACOTTA)
        add_text(s, Inches(4.5), y, Inches(8.5), Inches(0.4),
                 val, size=12, color=INK)
        y += Inches(0.6)
    slide_footer(s, 12, total)

    # ─── Slide 13 : À ne pas oublier / éviter ───
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    slide_header(s, "À mentionner / à éviter",
                 subtitle="Réflexes en démo")

    # À mentionner
    add_text(s, Inches(0.5), Inches(1.4), Inches(6), Inches(0.4),
             "✓  À mentionner", size=14, bold=True, color=SUCCESS)
    add_bullets(s, Inches(0.7), Inches(1.85), Inches(5.8), Inches(4), [
        "Données restent au Canada (Loi 25)",
        "STOP / OPT-OUT respecté by design (Loi 25 + CASL)",
        "Tu gardes la main en 1 clic",
        "Tu vois en temps réel ce que dit l'IA (pas de boîte noire)",
        "Architecture multi-tenant prête (scale possible)",
    ], size=12, color=INK)

    # À éviter
    add_text(s, Inches(6.8), Inches(1.4), Inches(6), Inches(0.4),
             "✗  À éviter", size=14, bold=True, color=DANGER)
    add_bullets(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(4), [
        "Mentionner voice cloning (hors scope MVP)",
        "Mentionner app mobile native (hors scope)",
        "Annoncer un dashboard analytique multi-courtier",
        "Nommer Claude / OpenAI / n8n (sauf si tech demande)",
        "Donner un prix précis sans valider en interne d'abord",
    ], size=12, color=INK)

    # Post-démo
    add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9), SAND)
    add_text(s, Inches(0.7), Inches(6.12), Inches(12), Inches(0.3),
             "Post-démo (5 min après le call)", size=12, bold=True, color=NAVY)
    add_text(s, Inches(0.7), Inches(6.44), Inches(12), Inches(0.45),
             "Envoyer email récap : Loom URL + dashboard de test (?courtier=...) + noter prochaine étape",
             size=11, color=INK)
    slide_footer(s, 13, total)

    out = Path(__file__).parent / "NextMove_Demo_Storyboard_2026-05-13.pptx"
    prs.save(out)
    print(f"✓ Saved: {out}")


if __name__ == "__main__":
    main()
