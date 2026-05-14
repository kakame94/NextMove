"""Klaris — Mini-deck GTM (FR). Generator python-pptx.

Output: Klaris_GTM_FR.pptx (9 slides, 16:9).
Audience: equipe interne (Dennis/Eliot/Walkens/Seydou), partenaires
proches, sync trimestriel.

Source: etude de marche v3 §4 + §7 + apprentissages WhatsApp.
Imports helpers + tokens marque depuis build_klaris_deck.py.
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from build_klaris_deck import (
    TERRACOTTA, TERRACOTTA_STRONG, TERRACOTTA_SOFT,
    CREAM, CREAM_DEEP, CHARCOAL, CHARCOAL_60, WHITE, BORDER,
    SUCCESS, WARN_AMBER, DESTRUCTIVE, INFO,
    BODY_FONT, MONO_FONT, SLIDE_W_IN, SLIDE_H_IN,
    add_bg, add_text, add_eyebrow,
    add_h1, add_card,
    _set_shape_alpha,
    add_grid_pattern, add_circle, add_mascot,
    add_footer, set_notes,
)

from build_klaris_investor_deck import (
    add_section_header, add_kpi_card, add_table_row, add_bullet_line,
)

TOTAL_SLIDES = 11


def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_grid_pattern(slide, Inches(0), Inches(0),
                     Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
                     step_in=0.45, alpha=0.30)

    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(8.4), Inches(1.5),
                                  Inches(4.5), Inches(4.5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = TERRACOTTA
    _set_shape_alpha(glow, 0.10)
    glow.line.fill.background()
    add_mascot(slide, Inches(10.65), Inches(3.75), Inches(2.6))

    add_text(slide, "K L A R I S", Inches(0.7), Inches(0.7),
             Inches(3), Pt(28), font=MONO_FONT, size=16, bold=True, color=CHARCOAL)
    add_text(slide, "GO-TO-MARKET · Q3 2026 · v1.0",
             Inches(0.7), Inches(1.15), Inches(6), Pt(18),
             font=MONO_FONT, size=10, color=CHARCOAL_60)

    add_text(slide, "Bootstrap. Quebec. 100 clients M12.",
             Inches(0.7), Inches(2.8), Inches(8.2), Inches(1.2),
             font=BODY_FONT, size=48, bold=True, color=CHARCOAL,
             line_spacing=1.05)
    add_text(slide, "Action Layer pour courtier solo bilingue.",
             Inches(0.7), Inches(3.95), Inches(8.2), Inches(1.0),
             font=BODY_FONT, size=26, bold=True, color=TERRACOTTA,
             line_spacing=1.05)
    add_text(slide,
             "Plan d'execution 24 mois — 4 canaux acquisition · 5 partenariats fondateurs · "
             "budget Y1 ~ 60 K$ net · jalon M12 = 150 K$ ARR.",
             Inches(0.7), Inches(5.1), Inches(8.5), Inches(1.4),
             font=BODY_FONT, size=14, color=CHARCOAL_60, line_spacing=1.45)

    add_text(slide,
             "Sync hebdo · revue trimestrielle · proprietaires : Dennis · Eliot · Walkens · Seydou",
             Inches(0.7), Inches(6.85), Inches(10), Pt(20),
             font=MONO_FONT, size=10, color=CHARCOAL_60)

    set_notes(slide, "Cover deck GTM interne. Document de reference equipe.")


def slide_02_cadre(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Cadre strategique",
                       "6 decisions binaires deja prises.")

    decisions = [
        ("Financement", "Bootstrap", "Pas de levee avant M18 + PMF prouve"),
        ("Geo Phase 1", "RMR Montreal", "Laval + Rive-Sud + Rive-Nord · M1-M12"),
        ("ICP prioritaire", "Solo bilingue", "2-7 ans exp · 100-500 K$ GCI · 12-15+ tx/an"),
        ("Positionnement", "Action Layer", "Pas CRM · interface FUB/Centiva via API"),
        ("Killer message", "« Klaris agit »", "« Votre CRM stocke passivement »"),
        ("Jalon M12", "100 clients · 150 K$ ARR", "NRR > 100 % · churn < 3 %/mois"),
    ]
    cw = Inches(6.0)
    ch = Inches(1.4)
    gap_x = Inches(0.20)
    gap_y = Inches(0.18)
    start_x = Inches(0.7)
    start_y = Inches(2.4)
    for i, (axis, val, sub) in enumerate(decisions):
        row = i // 2
        col = i % 2
        x = start_x + col * (cw + gap_x)
        y = start_y + row * (ch + gap_y)
        add_card(slide, x, y, cw, ch, fill=WHITE, border=BORDER, radius=0.10)
        add_text(slide, axis.upper(), x + Inches(0.30), y + Inches(0.18),
                 Inches(2.5), Inches(0.3),
                 font=MONO_FONT, size=9, bold=True, color=CHARCOAL_60)
        add_text(slide, val, x + Inches(0.30), y + Inches(0.48),
                 cw - Inches(0.6), Inches(0.5),
                 font=BODY_FONT, size=18, bold=True, color=TERRACOTTA)
        add_text(slide, sub, x + Inches(0.30), y + Inches(0.96),
                 cw - Inches(0.6), Inches(0.4),
                 font=BODY_FONT, size=11, color=CHARCOAL_60, line_spacing=1.3)

    add_footer(slide, 2, TOTAL_SLIDES)
    set_notes(slide,
              "Decisions binaires deja figees. Tout ce qui suit en decoule. "
              "Pas a rediscuter chaque semaine.")


def slide_03_persona(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Persona prioritaire",
                       "Le courtier solo bilingue — saturation capital temps.")

    # Left: profile card
    card_x = Inches(0.7)
    card_y = Inches(2.3)
    cw = Inches(6.0)
    ch = Inches(4.7)
    add_card(slide, card_x, card_y, cw, ch,
             fill=WHITE, border=BORDER, radius=0.12)

    # Avatar
    add_circle(slide, card_x + Inches(1.0), card_y + Inches(1.0),
               Inches(1.4), fill=TERRACOTTA_SOFT)
    add_text(slide, "J", card_x + Inches(0.55), card_y + Inches(0.55),
             Inches(0.9), Inches(0.9),
             font=BODY_FONT, size=44, bold=True, color=TERRACOTTA_STRONG,
             align=PP_ALIGN.CENTER)

    add_text(slide, "Joanel (archetype)",
             card_x + Inches(2.2), card_y + Inches(0.45),
             Inches(3.5), Inches(0.5),
             font=BODY_FONT, size=22, bold=True, color=CHARCOAL)
    add_text(slide, "Courtier independant · RMR Montreal · bilingue FR/EN",
             card_x + Inches(2.2), card_y + Inches(1.0),
             Inches(3.5), Inches(0.4),
             font=MONO_FONT, size=11, color=TERRACOTTA)

    fields = [
        ("Experience", "2-7 ans actif · OACIQ a jour"),
        ("Volume", "12-15+ transactions/an"),
        ("GCI", "100 K$ - 500 K$ · mid-market"),
        ("Tech stack", "Centris + iPhone + 1-2 outils SaaS"),
        ("Douleur #1", "Perdre lead car en visite quand il appelle"),
        ("Budget tech mensuel", "150-300 $ discretionnaires"),
        ("Decision logiciel", "Solo · pas comite franchise"),
    ]
    for i, (lbl, val) in enumerate(fields):
        fy = card_y + Inches(2.0) + i * Inches(0.36)
        add_text(slide, lbl, card_x + Inches(0.35), fy,
                 Inches(1.9), Inches(0.32),
                 font=MONO_FONT, size=10, bold=True, color=CHARCOAL_60)
        add_text(slide, val, card_x + Inches(2.25), fy,
                 cw - Inches(2.5), Inches(0.32),
                 font=BODY_FONT, size=11, color=CHARCOAL, line_spacing=1.25)

    # Right: anti-persona + WTP
    rx = Inches(7.0)
    ry = card_y
    rw = Inches(5.6)

    add_card(slide, rx, ry, rw, Inches(2.0),
             fill=DESTRUCTIVE, border=None, radius=0.10, shadow=False)
    add_text(slide, "QUI ON IGNORE", rx + Inches(0.3), ry + Inches(0.2),
             rw - Inches(0.6), Inches(0.3),
             font=MONO_FONT, size=10, bold=True, color=WHITE)
    anti = [
        "Junior < 2 ans · GCI < 40 K$ · churn eleve",
        "Top-team RE/MAX (deja captif Centiva)",
        "Sieges sociaux franchises (cycle 18 mois)",
        "Courtier desengage / temps partiel",
    ]
    for i, a in enumerate(anti):
        add_text(slide, f"• {a}", rx + Inches(0.3), ry + Inches(0.55) + i * Inches(0.32),
                 rw - Inches(0.6), Inches(0.3),
                 font=BODY_FONT, size=11, color=WHITE, line_spacing=1.3)

    # WTP triggers
    add_card(slide, rx, ry + Inches(2.2), rw, Inches(2.5),
             fill=SUCCESS, border=None, radius=0.10, shadow=False)
    add_text(slide, "DECLENCHEURS WTP", rx + Inches(0.3), ry + Inches(2.4),
             rw - Inches(0.6), Inches(0.3),
             font=MONO_FONT, size=10, bold=True, color=WHITE)
    triggers = [
        "Vient de perdre 1 transaction par delai reponse",
        "Cherche a remplacer adjointe humaine 2,5 K$/mois",
        "Loi 25 audit en cours — peur sanction CAI",
        "Compare a confrere bullpen qui a Klaris",
        "Saison printemps : flux entrant +40 %",
    ]
    for i, t in enumerate(triggers):
        add_text(slide, f"• {t}", rx + Inches(0.3), ry + Inches(2.75) + i * Inches(0.32),
                 rw - Inches(0.6), Inches(0.3),
                 font=BODY_FONT, size=11, color=WHITE, line_spacing=1.3)

    add_footer(slide, 3, TOTAL_SLIDES)
    set_notes(slide,
              "Persona Joanel = beta-testeur reel + archetype etude v3 §3. "
              "Anti-persona = courtiers a fort risque churn ou bloque par franchise. "
              "Triggers WTP = moments achat dans cycle annuel.")


def slide_04_canaux(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "4 canaux acquisition",
                       "Ordre de priorite descendant. CAC blended cible 1 200 $.")

    channels = [
        ("01", "Meta Ads",
         "Moteur principal",
         "Reels Douleur/Solution · audience custom registre OACIQ",
         "4-6 K$/mois", "5-8 clients/mois M6+",
         TERRACOTTA),
        ("02", "RDV OACIQ + APCIQ",
         "Legitimite institutionnelle",
         "Kiosque Palais congres · demo live scan QR → SMS test",
         "12-15 K$/evt", "Brand · 100-200 emails qualif",
         WARN_AMBER),
        ("03", "Outbound directeurs franchises",
         "Bottom-up persona 4",
         "LinkedIn + ref Belma · pitch forfait equipe outil recrutement",
         "VP Ventes 100 %", "1-3 agences/an signees",
         INFO),
        ("04", "Referral courtier-courtier",
         "CAC le plus bas",
         "100 $/100 $ in-app · signature SMS desactivable",
         "5 K$/an payouts", "Ratio cible > 0,15",
         SUCCESS),
    ]
    cw = Inches(2.95)
    ch = Inches(4.5)
    gap = Inches(0.15)
    start_x = Inches(0.7)
    y = Inches(2.4)
    for i, (num, name, role, how, budget, kpi, col) in enumerate(channels):
        x = start_x + i * (cw + gap)
        add_card(slide, x, y, cw, ch, fill=WHITE, border=BORDER, radius=0.10)
        # Top stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        x, y, cw, Inches(0.45))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = col
        stripe.line.fill.background()
        add_text(slide, num, x + Inches(0.25), y + Inches(0.08),
                 Inches(1.0), Inches(0.35),
                 font=MONO_FONT, size=14, bold=True, color=WHITE)
        add_text(slide, name, x + Inches(0.25), y + Inches(0.65),
                 cw - Inches(0.5), Inches(0.5),
                 font=BODY_FONT, size=16, bold=True, color=CHARCOAL,
                 line_spacing=1.15)
        add_text(slide, role.upper(), x + Inches(0.25), y + Inches(1.15),
                 cw - Inches(0.5), Inches(0.35),
                 font=MONO_FONT, size=9, bold=True, color=col)
        add_text(slide, how, x + Inches(0.25), y + Inches(1.55),
                 cw - Inches(0.5), Inches(1.4),
                 font=BODY_FONT, size=11, color=CHARCOAL_60, line_spacing=1.4)
        # Budget + KPI
        add_text(slide, "BUDGET", x + Inches(0.25), y + Inches(3.05),
                 cw - Inches(0.5), Inches(0.25),
                 font=MONO_FONT, size=8, bold=True, color=CHARCOAL_60)
        add_text(slide, budget, x + Inches(0.25), y + Inches(3.30),
                 cw - Inches(0.5), Inches(0.35),
                 font=BODY_FONT, size=13, bold=True, color=CHARCOAL,
                 line_spacing=1.2)
        add_text(slide, "OUTPUT", x + Inches(0.25), y + Inches(3.75),
                 cw - Inches(0.5), Inches(0.25),
                 font=MONO_FONT, size=8, bold=True, color=CHARCOAL_60)
        add_text(slide, kpi, x + Inches(0.25), y + Inches(4.00),
                 cw - Inches(0.5), Inches(0.4),
                 font=BODY_FONT, size=11, color=CHARCOAL, line_spacing=1.25)

    add_footer(slide, 4, TOTAL_SLIDES)
    set_notes(slide,
              "Canal 01 Meta = predictible · canal 03 Outbound = lumpy/cycle long. "
              "Diversification volontaire. Canal 04 Referral construit pendant M3+ "
              "quand premiers clients ravis.")


def slide_05_partenariats(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "5 partenariats fondateurs Annee 1",
                       "Un cofondateur owner par cible. Pas de partenariat sans nom.")

    headers = ["#", "Cible", "Objectif critique", "Owner", "Echeance"]
    rows = [
        ["1", "Centris Inc. (Direction Innovation / API)",
         "Accord syndication API MLS officiel",
         "Dennis", "M6"],
        ["2", "RE/MAX Anjou — Maxime Belma",
         "Pilote 50+ courtiers · success story chiffree",
         "Eliot", "M3"],
        ["3", "APCIQ (Communications + Evenements)",
         "Commandite RDV OACIQ + bilans stats",
         "Seydou", "M6"],
        ["4", "College Immobilier Quebec",
         "Inclusion cursus UFC technologies (formation continue)",
         "Walkens", "M9"],
        ["5", "Cabinet juridique Loi 25 (Lavery / BLG / Stikeman)",
         "Co-redaction livre blanc — lead magnet",
         "Walkens", "M4"],
    ]
    col_widths = [Inches(0.6), Inches(3.6), Inches(4.4), Inches(1.4), Inches(1.5)]
    row_h = Inches(0.7)
    x_start = Inches(0.7)
    y_start = Inches(2.4)
    add_table_row(slide, headers, x_start, y_start, col_widths,
                  Inches(0.5), header=True)
    for i, cells in enumerate(rows):
        add_table_row(slide, cells, x_start,
                      y_start + Inches(0.5) + i * row_h,
                      col_widths, row_h)

    # Callout
    add_text(slide,
             "Regle d'or : si un partenariat n'a pas d'owner ET d'echeance, il n'existe pas. "
             "Update statut sync hebdo.",
             Inches(0.7), Inches(6.55), Inches(SLIDE_W_IN - 1.4), Inches(0.5),
             font=BODY_FONT, size=12, italic=True, bold=True, color=TERRACOTTA_STRONG,
             line_spacing=1.3)

    add_footer(slide, 5, TOTAL_SLIDES)
    set_notes(slide,
              "Etude v3 §7.3. Centris = pilier survie. Belma = preuve sociale Persona 4. "
              "Cabinet juridique = lead magnet conformite Loi 25.")


def slide_06_sequence(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Sequence 24 mois",
                       "5 jalons. Trigger Phase 2 = cash-flow positif QC.")

    # Timeline strip
    miles = [
        ("M0-M3", "Setup", "Conv. cofondateurs · Tech E&O · subventions · entretiens Roof.ai",
         CHARCOAL_60),
        ("M3-M6", "Activation", "Meta 3 K$/mois · 10-20 payants · Belma pilote · VP Ventes nomine",
         WARN_AMBER),
        ("M6", "PMF proxy", "50 clients · NRR > 100 % cohort M3 · churn < 3 %/mois",
         TERRACOTTA),
        ("M6-M12", "Scale", "Meta 5-6 K$/mois · RDV OACIQ kiosque · Centris signe ou fallback DDF",
         TERRACOTTA),
        ("M12", "PMF valide", "100 clients · 150 K$ ARR · 1ere agence 10+ courtiers",
         SUCCESS),
        ("M18", "Decision", "Bootstrap pur OU prep Serie seed si NRR > 110 %",
         INFO),
        ("M24", "Trigger ON", "400 clients · 600 K$ ARR · cash-flow positif · adapter TRESA",
         SUCCESS),
    ]
    y = Inches(2.4)
    for i, (jalon, name, sub, col) in enumerate(miles):
        row_y = y + i * Inches(0.55)
        # Jalon pill
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.7), row_y,
                                      Inches(1.4), Inches(0.42))
        pill.adjustments[0] = 0.4
        pill.fill.solid()
        pill.fill.fore_color.rgb = col
        pill.line.fill.background()
        add_text(slide, jalon, Inches(0.7), row_y + Inches(0.05),
                 Inches(1.4), Inches(0.35),
                 font=MONO_FONT, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        # Name
        add_text(slide, name, Inches(2.3), row_y + Inches(0.05),
                 Inches(2.0), Inches(0.4),
                 font=BODY_FONT, size=14, bold=True, color=CHARCOAL,
                 line_spacing=1.1)
        # Sub
        add_text(slide, sub, Inches(4.4), row_y + Inches(0.08),
                 Inches(8.3), Inches(0.4),
                 font=BODY_FONT, size=11, color=CHARCOAL_60, line_spacing=1.3)

    add_text(slide,
             "Verdict actuel : Phase 1 GO Bootstrap · Phase 2 GO conditionnel · Phase 3 NO-GO.",
             Inches(0.7), Inches(6.55), Inches(SLIDE_W_IN - 1.4), Inches(0.5),
             font=MONO_FONT, size=10, italic=True, color=CHARCOAL_60)

    add_footer(slide, 6, TOTAL_SLIDES)
    set_notes(slide,
              "Etude v3 §7.1. Pas d'expansion ON/BC tant que cash-flow QC pas positif. "
              "Decision M18 binaire : Bootstrap pur (defaut) vs Serie seed (si NRR > 110 % "
              "et pipeline sature).")


def slide_07_budget(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Budget marketing Y1",
                       "Net cash-out ~ 36-60 K$. Subventions empilees pour amortir.")

    headers = ["Poste", "Budget annuel", "Source / note"]
    rows = [
        (["Meta Ads", "36 K$ (3 K$/mo × 12)",
          "Augmenter a 6 K$/mo si CAC < 1 200 $"], None),
        (["RDV OACIQ kiosque + demo", "12-15 K$",
          "1 evenement majeur Y1 · Palais congres"], None),
        (["Outbound (VP Ventes externe)", "0-40 K$",
          "Si cofondateur 100 %, 0 cash"], None),
        (["Referral payouts", "5 K$",
          "50 referrals × 100 $"], None),
        (["Contenu / landing / videos temoignages", "8 K$",
          "Production interne + freelance"], None),
        (["Tech E&O assurance", "5-10 K$",
          "Northbridge / Intact Tech / Assur360 — quote Q1"], None),
        (["SUBVENTIONS NON-DILUTIVES", "-30 K$",
          "PSCE PME MTL · Impulsion MEI · BDC Tech"], SUCCESS),
        (["NET CASH-OUT MARKETING Y1", "~ 36-60 K$",
          "Soutenable bootstrap · 100 clients × 100 $/mo Y1"], TERRACOTTA_SOFT),
    ]
    col_widths = [Inches(4.0), Inches(3.5), Inches(4.4)]
    row_h = Inches(0.45)
    x_start = Inches(0.7)
    y_start = Inches(2.4)
    add_table_row(slide, headers, x_start, y_start, col_widths,
                  Inches(0.5), header=True)
    for i, (cells, bg) in enumerate(rows):
        add_table_row(slide, cells, x_start,
                      y_start + Inches(0.5) + i * row_h,
                      col_widths, row_h, bg=bg)

    add_footer(slide, 7, TOTAL_SLIDES)
    set_notes(slide,
              "Budget bootstrap-realiste. Subventions empilees critique. "
              "Si PSCE + Impulsion + BDC Tech obtenus, net cash-out tombe a "
              "20-30 K$ pour Y1.")


def slide_08_kpi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "KPI dashboard",
                       "North-Star + 5 metriques surveillees chaque lundi.")

    # North-star banner
    add_card(slide, Inches(0.7), Inches(2.1),
             Inches(SLIDE_W_IN - 1.4), Inches(1.1),
             fill=CHARCOAL, border=None, radius=0.10, shadow=False)
    add_text(slide, "NORTH-STAR METRIC",
             Inches(0.95), Inches(2.25), Inches(8), Inches(0.3),
             font=MONO_FONT, size=10, bold=True, color=TERRACOTTA_SOFT)
    add_text(slide,
             "Nombre de rencontres physiques qualifiees confirmees autonomement par Klaris / semaine.",
             Inches(0.95), Inches(2.55), Inches(SLIDE_W_IN - 2.0), Inches(0.5),
             font=BODY_FONT, size=15, color=CREAM, line_spacing=1.3)

    # 5 KPI cards
    kpis = [
        ("MQL hebdo", "Cible 8-12", "Meta + LinkedIn + referral combinés", TERRACOTTA),
        ("CAC blended", "< 1 200 $", "Cash marketing ÷ nouveaux clients", WARN_AMBER),
        ("NRR cohort", "> 100 %", "Net revenue retention par cohorte mensuelle", SUCCESS),
        ("Churn logo", "< 3 %/mo", "Clients perdus / clients actifs · alerte > 5 %", INFO),
        ("Pipeline outbound", "Stages B2B", "Deals agence par etage funnel", DESTRUCTIVE),
    ]
    kw = Inches(2.42)
    kh = Inches(2.5)
    gap = Inches(0.10)
    start_x = Inches(0.7)
    y = Inches(3.55)
    for i, (lbl, val, sub, col) in enumerate(kpis):
        x = start_x + i * (kw + gap)
        add_kpi_card(slide, x, y, kw, kh, lbl, val, sub, value_color=col)

    # Rituel
    add_text(slide,
             "Rituel : sync hebdo lundi 10h · revue trim Q1/Q2/Q3/Q4 · escalation > 5 % churn.",
             Inches(0.7), Inches(6.45), Inches(SLIDE_W_IN - 1.4), Inches(0.5),
             font=MONO_FONT, size=10, italic=True, color=CHARCOAL_60)

    add_footer(slide, 8, TOTAL_SLIDES)
    set_notes(slide,
              "North-Star anti-vanity = pas downloads / visites. RDV physique = "
              "preuve ROI. 5 KPI hebdo permettent detection precoce derive.")


def slide_09_actions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_grid_pattern(slide, Inches(0), Inches(0),
                     Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
                     step_in=0.45, alpha=0.25)

    add_section_header(slide, "Actions semaine 1",
                       "5 decisions binaires. Tout le reste en decoule.")

    actions = [
        ("1.", "VP Ventes nomine",
         "1 cofondateur 100 % dedie pipeline agence — Dennis/Eliot/Walkens/Seydou."),
        ("2.", "Tech E&O quote",
         "Appeler Northbridge + Intact Tech + Assur360 — budget 5-10 K$/an."),
        ("3.", "Convention cofondateurs",
         "Cap table + vesting 4 ans + cliff 1 an — chez notaire ou Clerky."),
        ("4.", "5 entretiens Roof.ai ex-clients",
         "Recherche primaire critique etude v3 — churn driver concurrent local."),
        ("5.", "Inscription RDV OACIQ 2026",
         "Calendrier evenement + budget kiosque + slot demo live."),
    ]
    cw = Inches(SLIDE_W_IN - 1.4)
    ch = Inches(0.7)
    gap_y = Inches(0.12)
    start_x = Inches(0.7)
    start_y = Inches(2.3)
    for i, (num, title, sub) in enumerate(actions):
        y = start_y + i * (ch + gap_y)
        add_card(slide, start_x, y, cw, ch, fill=WHITE, border=BORDER, radius=0.10)
        add_text(slide, num, start_x + Inches(0.3), y + Inches(0.12),
                 Inches(0.7), Inches(0.5),
                 font=MONO_FONT, size=22, bold=True, color=TERRACOTTA)
        add_text(slide, title, start_x + Inches(1.1), y + Inches(0.07),
                 Inches(4.5), Inches(0.35),
                 font=BODY_FONT, size=15, bold=True, color=CHARCOAL,
                 line_spacing=1.15)
        add_text(slide, sub, start_x + Inches(1.1), y + Inches(0.36),
                 cw - Inches(1.4), Inches(0.32),
                 font=BODY_FONT, size=11, color=CHARCOAL_60, line_spacing=1.3)

    # Final call
    add_text(slide,
             "« Pas de levee. Pas de Centiva head-on. Bootstrap pur jusqu'a 100 clients prouves. »",
             Inches(0.7), Inches(6.65), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
             font=BODY_FONT, size=13, italic=True, bold=True, color=TERRACOTTA_STRONG,
             align=PP_ALIGN.CENTER, line_spacing=1.3)

    add_footer(slide, 9, TOTAL_SLIDES)
    set_notes(slide,
              "Close GTM deck. 5 actions tangibles cette semaine. Sans elles, "
              "rien d'autre n'avance. Owner assigne par cofondateur, revue "
              "vendredi prochain.")


def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slide_01_cover(prs)
    slide_02_cadre(prs)
    slide_03_persona(prs)
    slide_04_canaux(prs)
    slide_05_partenariats(prs)
    slide_06_sequence(prs)
    slide_07_budget(prs)
    slide_08_kpi(prs)
    slide_09_actions(prs)

    from klaris_annexes import slide_annexe_acronymes, slide_annexe_metriques
    slide_annexe_acronymes(prs, 10, TOTAL_SLIDES)
    slide_annexe_metriques(prs, 11, TOTAL_SLIDES)

    out = Path(__file__).parent / "Klaris_GTM_FR.pptx"
    prs.save(out)
    print(f"OK -> {out}")


if __name__ == "__main__":
    build()
