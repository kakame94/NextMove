"""Klaris — Deck avant-vente solo FR. Generator python-pptx.

Output: Klaris_AvantVente_Solo_FR.pptx (14 slides, 16:9).
DA: terracotta #ff5a36 + cream + Geist (alignee sur klaris-presentation.html).
Mockups iPhone/dashboard composes en shapes natives.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

ASSETS = Path(__file__).parent / "assets_pptx"
LOGO_MARK_PNG = str(ASSETS / "klaris_mark.png")          # mark "C" only (square)
LOGO_HORIZ_PNG = str(ASSETS / "klaris_logo.png")          # mark + wordmark Klaris (wide)
LOGO_STACKED_PNG = str(ASSETS / "klaris_stacked.png")     # mark + wordmark stacked
ROBOT_PNG = str(ASSETS / "klaris_robot.png")              # robot mascotte


# ----------------------------------------------------------------------------
# DA — palette + type (source: klaris_ios/lib/core/theme/klaris_colors.dart)
# ----------------------------------------------------------------------------

TERRACOTTA       = RGBColor(0xFF, 0x5A, 0x36)   # match site --primary (favicon literal)
TERRACOTTA_STRONG= RGBColor(0xCC, 0x48, 0x28)   # site --primary-strong (hover)
TERRACOTTA_SOFT  = RGBColor(0xFF, 0xE2, 0xD6)   # ~10% terracotta on cream
CREAM            = RGBColor(0xFA, 0xF9, 0xF5)
CREAM_DEEP       = RGBColor(0xEE, 0xEA, 0xE2)   # muted
CHARCOAL         = RGBColor(0x2A, 0x25, 0x21)
CHARCOAL_60      = RGBColor(0x7A, 0x70, 0x68)   # mutedFg
WHITE            = RGBColor(0xFF, 0xFF, 0xFF)
BORDER           = RGBColor(0xE2, 0xDD, 0xD3)
SUCCESS          = RGBColor(0x50, 0xA0, 0x4B)
WARN_AMBER       = RGBColor(0xE0, 0xA8, 0x36)
DESTRUCTIVE      = RGBColor(0xD6, 0x4C, 0x2E)
INFO             = RGBColor(0x4F, 0x8F, 0xE0)
INCOMING_BUBBLE  = RGBColor(0xEF, 0xED, 0xE6)

HEAT = [
    RGBColor(0x6F, 0xA1, 0xD8),  # froid
    RGBColor(0x8A, 0xB4, 0xD0),  # tiede
    RGBColor(0xE0, 0xA8, 0x36),  # chaud (ambre)
    RGBColor(0xE0, 0x7A, 0x45),  # tres chaud
    RGBColor(0xD6, 0x4C, 0x2E),  # brulant
]

BODY_FONT = "Geist"
MONO_FONT = "Geist Mono"

# 16:9
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------

def add_bg(slide, color=CREAM):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def _set_text(tf, text, *, font=BODY_FONT, size=14, bold=False, italic=False,
              color=CHARCOAL, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.15):
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run() if p.runs == [] else p.runs[0]
    if not p.runs:
        run = p.add_run()
    else:
        # clear default
        run = p.runs[0]
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color


def add_text(slide, text, x, y, w, h, *, font=BODY_FONT, size=14, bold=False,
             italic=False, color=CHARCOAL, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return tb


def add_multitext(slide, lines, x, y, w, h, *, default=None, anchor=MSO_ANCHOR.TOP):
    """lines: list of dicts {text, font, size, bold, italic, color, align, space_before}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    d = default or {}
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        spec = {**d, **line}
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        p.line_spacing = spec.get("line_spacing", 1.2)
        if "space_before" in spec:
            p.space_before = Pt(spec["space_before"])
        if "space_after" in spec:
            p.space_after = Pt(spec["space_after"])
        run = p.add_run()
        run.text = spec["text"]
        f = run.font
        f.name = spec.get("font", BODY_FONT)
        f.size = Pt(spec.get("size", 14))
        f.bold = spec.get("bold", False)
        f.italic = spec.get("italic", False)
        f.color.rgb = spec.get("color", CHARCOAL)
    return tb


def add_eyebrow(slide, text, x, y, w=Inches(8), color=TERRACOTTA):
    """Pill badge style (motif site Klaris): bg primary-soft + dot + text primary."""
    # Estimate text width: ~7pt per char + dot + padding
    label = text.upper()
    pill_w = Inches(0.40 + len(label) * 0.085)
    pill_h = Inches(0.32)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, pill_w, pill_h)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = TERRACOTTA
    _set_shape_alpha(pill, 0.10)
    pill.line.fill.background()
    # Dot
    dot_d = Inches(0.10)
    dot_cx = x + Inches(0.18)
    dot_cy = y + pill_h / 2
    add_circle(slide, dot_cx, dot_cy, dot_d, fill=color)
    # Text
    add_text(slide, label,
             x + Inches(0.30), y + Inches(0.06),
             pill_w - Inches(0.40), pill_h - Inches(0.10),
             font=MONO_FONT, size=9, bold=True, color=color, line_spacing=1.0)


def add_logo_mark(slide, x, y, size):
    """Embed Klaris mark PNG (C with chat tail + 3 dots)."""
    return slide.shapes.add_picture(LOGO_MARK_PNG, x, y, height=size, width=size)


def add_logo_horizontal(slide, x, y, height):
    """Embed Klaris horizontal logo (mark + wordmark Klaris)."""
    return slide.shapes.add_picture(LOGO_HORIZ_PNG, x, y, height=height)


def add_logo_stacked(slide, x, y, height):
    return slide.shapes.add_picture(LOGO_STACKED_PNG, x, y, height=height)


def add_robot(slide, x, y, height):
    """Embed robot mascotte PNG."""
    return slide.shapes.add_picture(ROBOT_PNG, x, y, height=height)


def add_h2(slide, text, x, y, w, *, color=CHARCOAL, size=24, align=PP_ALIGN.LEFT):
    return add_text(slide, text, x, y, w, Pt(size * 2),
                    font=BODY_FONT, size=size, bold=True, color=color,
                    align=align, line_spacing=1.15)


def add_h1(slide, text, x, y, w, *, color=CHARCOAL, size=36, align=PP_ALIGN.LEFT):
    return add_text(slide, text, x, y, w, Pt(size * 2),
                    font=BODY_FONT, size=size, bold=True, color=color,
                    align=align, line_spacing=1.1)


def add_body(slide, text, x, y, w, h, *, color=CHARCOAL, size=14, align=PP_ALIGN.LEFT):
    return add_text(slide, text, x, y, w, h,
                    font=BODY_FONT, size=size, color=color, align=align,
                    line_spacing=1.35)


def add_card(slide, x, y, w, h, *, fill=WHITE, border=None, radius=0.08, shadow=True):
    if shadow:
        # ombre simulee: rect charcoal opacity ~8% offset 3pt
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x + Pt(3), y + Pt(4), w, h)
        sh.adjustments[0] = radius
        sh.fill.solid()
        sh.fill.fore_color.rgb = CHARCOAL
        _set_shape_alpha(sh, 0.08)
        sh.line.fill.background()
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = radius
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border is not None:
        card.line.color.rgb = border
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    return card


def _set_shape_alpha(shape, alpha):
    """Apply alpha (0-1) on solid fill via raw XML."""
    sp = shape.fill._xPr  # solidFill parent
    solid = sp.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    # remove existing alpha
    for a in srgb.findall(qn("a:alpha")):
        srgb.remove(a)
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", str(int(alpha * 100000)))


def add_grid_pattern(slide, x, y, w, h, *, line_color=BORDER, alpha=0.55,
                     step_in=0.55, dotted=False, on_dark=False):
    """Subtle blueprint grid pattern — match web `.blueprint-bg::before`.

    Drawn as a dense set of thin lines on top of [x,y,w,h] region.
    Use behind hero / archi sections for the sci-fi engineering feel.
    """
    color = WHITE if on_dark else line_color
    base_alpha = alpha if not on_dark else 0.10

    # Vertical lines
    cur = x + Inches(step_in)
    end_x = x + w
    while cur < end_x:
        ln = slide.shapes.add_connector(1, cur, y, cur, y + h)
        ln.line.color.rgb = color
        ln.line.width = Pt(0.4)
        _set_line_alpha(ln, base_alpha)
        cur += Inches(step_in)

    # Horizontal lines
    cur = y + Inches(step_in)
    end_y = y + h
    while cur < end_y:
        ln = slide.shapes.add_connector(1, x, cur, x + w, cur)
        ln.line.color.rgb = color
        ln.line.width = Pt(0.4)
        _set_line_alpha(ln, base_alpha)
        cur += Inches(step_in)


def _set_line_alpha(shape, alpha):
    """Apply alpha on solid line via raw XML."""
    spPr = shape.line._get_or_add_ln()
    solid = spPr.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for a in srgb.findall(qn("a:alpha")):
        srgb.remove(a)
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", str(int(alpha * 100000)))


def add_circle(slide, cx, cy, d, *, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   cx - d / 2, cy - d / 2, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_heat_dots(slide, x, y, level=4, *, dot_d=Pt(11), gap=Pt(5)):
    """level 1-5: nb dots allumes."""
    shapes = []
    for i in range(5):
        col = HEAT[i]
        cx = x + dot_d / 2 + i * (dot_d + gap)
        cy = y + dot_d / 2
        s = add_circle(slide, cx, cy, dot_d, fill=col)
        if i >= level:
            _set_shape_alpha(s, 0.20)
        shapes.append(s)
    total_w = 5 * dot_d + 4 * gap
    return total_w


def add_mascot(slide, cx, cy, size, *, body=TERRACOTTA, dot=WHITE):
    # Cercle terracotta + 3 dots blancs alignes horizontal
    add_circle(slide, cx, cy, size, fill=body)
    dot_d = size * 0.13
    spacing = size * 0.22
    for i, dx in enumerate([-spacing, 0, spacing]):
        add_circle(slide, cx + dx, cy, dot_d, fill=dot)


def add_button(slide, x, y, w, h, label, *, primary=True):
    if primary:
        fill = TERRACOTTA
        border = None
        text_color = WHITE
    else:
        fill = WHITE
        border = TERRACOTTA
        text_color = TERRACOTTA
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    btn.adjustments[0] = 0.45  # tres arrondi
    btn.fill.solid()
    btn.fill.fore_color.rgb = fill
    if border is not None:
        btn.line.color.rgb = border
        btn.line.width = Pt(1.5)
    else:
        btn.line.fill.background()
    tf = btn.text_frame
    tf.margin_left = tf.margin_right = Pt(8)
    tf.margin_top = tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    f = run.font
    f.name = BODY_FONT
    f.size = Pt(11)
    f.bold = True
    f.color.rgb = text_color
    return btn


def add_footer(slide, page_num, total=14):
    # wordmark gauche
    add_text(slide, "K L A R I S", Inches(0.5), Inches(SLIDE_H_IN - 0.42),
             Inches(2), Pt(14),
             font=MONO_FONT, size=8, bold=True, color=CHARCOAL_60)
    # page count droite
    add_text(slide, f"{page_num:02d} / {total:02d}",
             Inches(SLIDE_W_IN - 1.5), Inches(SLIDE_H_IN - 0.42),
             Inches(1), Pt(14),
             font=MONO_FONT, size=8, bold=True, color=CHARCOAL_60,
             align=PP_ALIGN.RIGHT)


def set_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


# ----------------------------------------------------------------------------
# Slides
# ----------------------------------------------------------------------------

def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)

    # Blueprint grid (signature web — engineering / sci-fi feel)
    add_grid_pattern(slide, Inches(0), Inches(0),
                     Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
                     step_in=0.45, alpha=0.35)

    # Soft radial glow behind mascot — emulates web `.hero-visual::before`
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(7.6), Inches(1.0),
                                  Inches(6.0), Inches(6.0))
    glow.fill.solid()
    glow.fill.fore_color.rgb = TERRACOTTA
    _set_shape_alpha(glow, 0.06)
    glow.line.fill.background()

    # Robot mascotte a droite (hero visual)
    add_robot(slide, Inches(8.5), Inches(1.4), Inches(5.0))

    # Logo horizontal Klaris (mark + wordmark) en haut-gauche
    add_logo_horizontal(slide, Inches(0.6), Inches(0.6), Inches(0.85))

    # Hero H1
    add_text(slide,
             "L'adjointe IA des",
             Inches(0.6), Inches(2.6),
             Inches(7.5), Pt(70),
             font=BODY_FONT, size=52, bold=True, color=CHARCOAL,
             line_spacing=1.05)
    add_text(slide,
             "courtiers immo.",
             Inches(0.6), Inches(3.5),
             Inches(7.5), Pt(70),
             font=BODY_FONT, size=52, bold=True, color=TERRACOTTA,
             line_spacing=1.05)

    # Sub
    add_text(slide,
             "Elle qualifie tes prospects par SMS. Toi, tu vends.",
             Inches(0.6), Inches(4.7), Inches(8), Pt(36),
             font=BODY_FONT, size=18, color=CHARCOAL_60, line_spacing=1.4)

    # Eyebrow bas-gauche (style web: "Live · ... · v2.1")
    add_eyebrow(slide, "Live  ·  pilote courtiers Grand Montreal  ·  v2.1",
                Inches(0.6), Inches(SLIDE_H_IN - 0.95))

    # Credit bas-droite — domain klarisapp.ai
    add_text(slide, "klarisapp.ai  ·  Une marque Next Move",
             Inches(SLIDE_W_IN - 5.0), Inches(SLIDE_H_IN - 0.55),
             Inches(4.4), Pt(14),
             font=MONO_FONT, size=9, bold=True, color=CHARCOAL_60,
             align=PP_ALIGN.RIGHT)

    set_notes(slide,
              "Bonjour, je vous presente Klaris. C'est l'adjointe IA des courtiers "
              "immobiliers. Elle qualifie tes prospects par SMS — pendant que toi, "
              "tu vends. (Pause apres la tagline. Laisse landing.)")
    return slide


def slide_02_probleme(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "Le probleme", Inches(0.6), Inches(0.6))
    add_h2(slide, "L'admin tue les courtiers. Pas le marche.",
           Inches(0.6), Inches(1.0), Inches(11), size=28)

    cards = [
        ("30-60 min", "pour rediger une offre d'achat a la main"),
        ("+24h", "de delai moyen avant de repondre a un nouveau lead"),
        ("80%", "des CRMs de franchise sont vides — saisie manuelle = adoption nulle"),
    ]
    card_w = Inches(3.8)
    card_h = Inches(3.2)
    gap = Inches(0.4)
    total_w = Emu(3 * card_w + 2 * gap)
    start_x = Inches((SLIDE_W_IN - total_w.inches) / 2)
    y = Inches(2.55)

    for i, (big, sub) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, y, card_w, card_h)
        # Top accent bar destructive (motif "pain card" du site Klaris)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        x + Inches(0.18), y,
                                        card_w - Inches(0.36), Pt(3.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = DESTRUCTIVE
        accent.line.fill.background()
        # gros chiffre (size 36 evite wrap "30-60 min")
        add_text(slide, big, x + Inches(0.3), y + Inches(0.6),
                 card_w - Inches(0.6), Pt(60),
                 font=BODY_FONT, size=36, bold=True, color=TERRACOTTA,
                 align=PP_ALIGN.LEFT, line_spacing=1.05)
        # sous-texte (descend pour eviter overlap)
        add_text(slide, sub, x + Inches(0.3), y + Inches(1.8),
                 card_w - Inches(0.6), Inches(1.4),
                 font=BODY_FONT, size=13, color=CHARCOAL, line_spacing=1.4)

    add_text(slide,
             "Source : 4 entretiens approfondis, courtiers terrain QC, mars 2026",
             Inches(0.6), Inches(6.0), Inches(11), Pt(20),
             font=MONO_FONT, size=9, italic=True, color=CHARCOAL_60)

    add_footer(slide, 2)
    set_notes(slide,
              "Avant de parler de Klaris, parlons du probleme. On a interviewe 4 "
              "courtiers terrain en mars. Tous les quatre, sans exception, nous "
              "ont dit la meme chose : leur probleme n1 n'est pas le marche, c'est "
              "l'admin. Ce n'est pas un probleme d'effort. C'est un probleme de "
              "systeme.")
    return slide


def slide_03_cout_inaction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "Cout de l'inaction", Inches(0.6), Inches(0.6))
    add_h2(slide, "Le courtier perd ce qu'il fait le mieux.",
           Inches(0.6), Inches(1.0), Inches(11), size=28)

    cards = [
        ("Temps perdu", "2h+/jour en admin = 10h+/semaine = 1 transaction/mois ratee", "icon_clock"),
        ("Prospects perdus", "Un lead chaud sans reponse en 1h se refroidit, parle a la concurrence", "icon_thermo"),
        ("Plafond impose", "Sans systemes, impossible de scaler de 3 a 10 transactions/mois", "icon_chart"),
    ]
    card_w = Inches(3.8)
    card_h = Inches(3.7)
    gap = Inches(0.4)
    total_w = Emu(3 * card_w + 2 * gap)
    start_x = Inches((SLIDE_W_IN - total_w.inches) / 2)
    y = Inches(2.4)

    for i, (title, body, icon) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, y, card_w, card_h)
        # icone (cercle terracotta soft + symbole charcoal)
        ic_cx = x + Inches(0.6)
        ic_cy = y + Inches(0.65)
        add_circle(slide, ic_cx, ic_cy, Inches(0.55), fill=TERRACOTTA_SOFT)
        # symbole simple lettre dans le cercle
        sym = {"icon_clock": "T", "icon_thermo": "P", "icon_chart": "X"}[icon]
        add_text(slide, sym,
                 ic_cx - Inches(0.4), ic_cy - Inches(0.2),
                 Inches(0.8), Inches(0.4),
                 font=BODY_FONT, size=18, bold=True, color=TERRACOTTA,
                 align=PP_ALIGN.CENTER)
        # titre
        add_text(slide, title, x + Inches(0.4), y + Inches(1.4),
                 card_w - Inches(0.8), Pt(28),
                 font=BODY_FONT, size=18, bold=True, color=CHARCOAL)
        # body
        add_body(slide, body, x + Inches(0.4), y + Inches(2.0),
                 card_w - Inches(0.8), Inches(1.5),
                 size=13, color=CHARCOAL_60)

    add_footer(slide, 3)
    set_notes(slide,
              "Concretement, ca coute quoi ? Premier cout: le temps. Deuxieme: les "
              "prospects. Troisieme: le plafond. Sans systeme, pas moyen de scaler.")
    return slide


def slide_04_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "La solution", Inches(0.6), Inches(0.6))

    # Robot geant a gauche
    add_robot(slide, Inches(0.8), Inches(1.6), Inches(5.4))

    # Bloc texte a droite
    rx = Inches(7.0)
    add_logo_horizontal(slide, rx, Inches(2.0), Inches(0.7))

    add_text(slide,
             "Klaris qualifie tes prospects par SMS.",
             rx, Inches(3.05), Inches(6.0), Pt(50),
             font=BODY_FONT, size=24, bold=True, color=CHARCOAL,
             line_spacing=1.3)
    add_text(slide,
             "Pendant que tu vends.",
             rx, Inches(3.95), Inches(6.0), Pt(40),
             font=BODY_FONT, size=24, bold=True, color=TERRACOTTA,
             line_spacing=1.3)

    add_text(slide,
             "Un seul outil. Un seul abonnement. Conformite OACIQ par defaut. Bilingue FR / EN.",
             rx, Inches(4.85), Inches(6.0), Inches(1.0),
             font=BODY_FONT, size=14, color=CHARCOAL_60, line_spacing=1.45)

    # Tagline italic dans pill terracotta-soft
    tagline_w = Inches(2.8)
    tagline_h = Inches(0.45)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  rx, Inches(6.0), tagline_w, tagline_h)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = TERRACOTTA
    _set_shape_alpha(pill, 0.10)
    pill.line.fill.background()
    add_text(slide, "« Je suis son assistante. »",
             rx, Inches(6.06), tagline_w, Pt(28),
             font=BODY_FONT, size=12, italic=True, bold=True, color=TERRACOTTA,
             align=PP_ALIGN.CENTER)

    add_footer(slide, 4)
    set_notes(slide,
              "Notre solution s'appelle Klaris. C'est une assistante IA specialisee "
              "pour les courtiers immo. Trois principes : un seul outil, un seul "
              "abonnement, conformite OACIQ par defaut. Et — detail important — "
              "Klaris se presente toujours comme l'assistante du courtier. Jamais "
              "comme un robot.")
    return slide


def slide_05_comment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "Comment ca marche", Inches(0.6), Inches(0.6))
    add_h2(slide, "SMS arrive. Klaris qualifie. Tu fermes.",
           Inches(0.6), Inches(1.0), Inches(11), size=28)

    steps = [
        ("1", "Le prospect texte", "ton numero habituel."),
        ("2", "Klaris pose 5 questions cles", "et resume — ton chaleureux Quebecois."),
        ("3", "Le lead chaud apparait", "dans ton dashboard avec un score de temperature."),
    ]
    circle_d = Inches(1.0)
    col_w = Inches(3.6)
    gap = Inches(0.5)
    total_w = Emu(3 * col_w + 2 * gap)
    start_x = Inches((SLIDE_W_IN - total_w.inches) / 2)
    y = Inches(3.0)

    # Connector ligne horizontale arriere-plan (gradient terracotta-soft -> terracotta)
    line_y = y
    line_x1 = start_x + col_w / 2
    line_x2 = start_x + 2 * col_w + 2 * gap + col_w / 2
    connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       line_x1, line_y - Pt(1),
                                       line_x2 - line_x1, Pt(2))
    connector.fill.solid()
    connector.fill.fore_color.rgb = TERRACOTTA
    _set_shape_alpha(connector, 0.30)
    connector.line.fill.background()

    for i, (num, title, body) in enumerate(steps):
        x = start_x + i * (col_w + gap)
        cx = x + col_w / 2
        # cercle blanc avec bordure terracotta (motif site Klaris)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        cx - circle_d / 2, y - circle_d / 2,
                                        circle_d, circle_d)
        circle.fill.solid()
        circle.fill.fore_color.rgb = WHITE
        circle.line.color.rgb = TERRACOTTA
        circle.line.width = Pt(2)
        # numero terracotta
        add_text(slide, num,
                 cx - Inches(0.5), y - Inches(0.32),
                 Inches(1.0), Inches(0.8),
                 font=BODY_FONT, size=28, bold=True, color=TERRACOTTA,
                 align=PP_ALIGN.CENTER)
        # titre
        add_text(slide, title,
                 x, y + Inches(0.85),
                 col_w, Pt(28),
                 font=BODY_FONT, size=16, bold=True, color=CHARCOAL,
                 align=PP_ALIGN.CENTER, line_spacing=1.25)
        # body
        add_body(slide, body, x, y + Inches(1.45),
                 col_w, Inches(1.5),
                 size=13, color=CHARCOAL_60, align=PP_ALIGN.CENTER)

    add_footer(slide, 5)
    set_notes(slide,
              "Trois etapes. Un : le prospect texte ton numero habituel. Pas de QR "
              "code, pas d'app a installer. Deux : Klaris pose 5 questions cles. "
              "Trois : le lead arrive dans ton dashboard avec un score de "
              "temperature. Tu rappelles quand le moment est bon.")
    return slide


def slide_06_demo_sms(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "Demo", Inches(0.6), Inches(0.6))
    add_h2(slide, "Une vraie conversation. Pas un formulaire.",
           Inches(0.6), Inches(1.0), Inches(7.5), size=26)

    # Bloc gauche — contexte
    left_x = Inches(0.6)
    left_y = Inches(2.4)
    add_text(slide,
             "Voici ce qui se passe pendant que tu es en visite.",
             left_x, left_y, Inches(7), Pt(30),
             font=BODY_FONT, size=15, color=CHARCOAL, line_spacing=1.4)

    bullets = [
        ("Le prospect texte ton numero a 19h47", left_y + Inches(0.7)),
        ("Klaris repond en moins de 10 secondes", left_y + Inches(1.4)),
        ("Tu recois la fiche prospect prete au prochain coup d'oeil", left_y + Inches(2.1)),
    ]
    for b, by in bullets:
        # Bullet dot
        add_circle(slide, left_x + Inches(0.12), by + Inches(0.15),
                   Inches(0.12), fill=TERRACOTTA)
        add_text(slide, b, left_x + Inches(0.4), by,
                 Inches(6.5), Pt(28),
                 font=BODY_FONT, size=14, color=CHARCOAL, line_spacing=1.4)

    add_text(slide, "C O N V E R S A T I O N   R E E L L E   ·   P I L O T E   T E R R A I N",
             left_x, left_y + Inches(3.0), Inches(7), Pt(20),
             font=MONO_FONT, size=9, bold=True, color=TERRACOTTA)

    # iPhone mockup a droite
    iphone_x = Inches(8.7)
    iphone_y = Inches(0.8)
    iphone_w = Inches(3.8)
    iphone_h = Inches(6.2)
    draw_iphone_sms(slide, iphone_x, iphone_y, iphone_w, iphone_h)

    add_footer(slide, 6)
    set_notes(slide,
              "Voila a quoi ressemble une vraie conversation. Une question a la "
              "fois. Pas de mur de texte. Ton naturel. C'est calibre sur les "
              "courtiers du terrain. Et pendant que tu visites une autre propriete, "
              "Klaris fait tout ca.")
    return slide


def draw_iphone_sms(slide, x, y, w, h):
    # Cadre exterieur (frame charcoal)
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    frame.adjustments[0] = 0.10
    frame.fill.solid()
    frame.fill.fore_color.rgb = CHARCOAL
    frame.line.fill.background()

    # Ecran interieur
    pad = Pt(6)
    sw = w - 2 * pad
    sh = h - 2 * pad
    sx = x + pad
    sy = y + pad
    screen = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, sy, sw, sh)
    screen.adjustments[0] = 0.085
    screen.fill.solid()
    screen.fill.fore_color.rgb = CREAM
    screen.line.fill.background()

    # Notch
    notch_w = Inches(1.2)
    notch_h = Inches(0.18)
    nx = x + (w - notch_w) / 2
    ny = y + Pt(8)
    notch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, nx, ny, notch_w, notch_h)
    notch.adjustments[0] = 0.5
    notch.fill.solid()
    notch.fill.fore_color.rgb = CHARCOAL
    notch.line.fill.background()

    # Header bar (terracotta)
    header_y = sy + Inches(0.3)
    header_h = Inches(0.55)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    sx, header_y, sw, header_h)
    header.fill.solid()
    header.fill.fore_color.rgb = TERRACOTTA
    header.line.fill.background()
    # Avatar mini avec mark Klaris
    avatar_d = Inches(0.36)
    add_circle(slide, sx + Inches(0.30), header_y + header_h / 2,
               avatar_d, fill=WHITE)
    slide.shapes.add_picture(LOGO_MARK_PNG,
                             sx + Inches(0.30) - avatar_d / 2 + Pt(2),
                             header_y + header_h / 2 - avatar_d / 2 + Pt(2),
                             height=avatar_d - Pt(4))
    # Nom
    add_text(slide, "Joanel — Klaris",
             sx + Inches(0.60), header_y + Inches(0.13),
             Inches(2.5), Inches(0.3),
             font=BODY_FONT, size=11, bold=True, color=WHITE)

    # Bubbles
    bubbles = [
        ("Bonsoir, je cherche un duplex a Verdun, budget 450K", "in",  "19:47"),
        ("Bonsoir ! Merci de contacter Joanel. Je suis son assistante. Acheter ou vendre ?", "out", "19:47"),
        ("Acheter", "in",  "19:48"),
        ("Super ! Dans quel secteur tu cherches ?", "out", "19:48"),
        ("Verdun ou Saint-Henri idealement", "in", "19:49"),
    ]
    bub_y = header_y + header_h + Inches(0.18)
    bub_max_w = sw - Inches(0.6)

    for text, kind, ts in bubbles:
        bw = min(bub_max_w, Inches(0.16) + Inches(0.085) * len(text))
        bw = min(bw, Inches(2.6))
        # estimate height by text length
        bh = Inches(0.28) + Inches(0.16) * max(1, len(text) // 25)
        if kind == "in":
            bx = sx + Inches(0.18)
            fill = INCOMING_BUBBLE
            txt_color = CHARCOAL
            ts_align = PP_ALIGN.LEFT
            ts_x = bx
        else:
            bx = sx + sw - bw - Inches(0.18)
            fill = TERRACOTTA
            txt_color = WHITE
            ts_align = PP_ALIGN.RIGHT
            ts_x = sx + sw - Inches(1.0) - Inches(0.18)
        bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        bx, bub_y, bw, bh)
        bubble.adjustments[0] = 0.30
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = fill
        bubble.line.fill.background()
        # texte dans la bubble
        tf = bubble.text_frame
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(4)
        tf.margin_bottom = Pt(4)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.2
        run = p.add_run()
        run.text = text
        run.font.name = BODY_FONT
        run.font.size = Pt(9)
        run.font.color.rgb = txt_color
        # timestamp
        add_text(slide, ts, ts_x, bub_y + bh + Pt(1),
                 Inches(1.0), Inches(0.16),
                 font=MONO_FONT, size=7, color=CHARCOAL_60, align=ts_align)
        bub_y = bub_y + bh + Inches(0.24)

    # Input bar bas
    input_h = Inches(0.4)
    input_y = sy + sh - input_h - Inches(0.15)
    input_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       sx + Inches(0.18), input_y,
                                       sw - Inches(0.36), input_h)
    input_bar.adjustments[0] = 0.5
    input_bar.fill.solid()
    input_bar.fill.fore_color.rgb = WHITE
    input_bar.line.color.rgb = BORDER
    input_bar.line.width = Pt(0.75)
    add_text(slide, "Message…",
             sx + Inches(0.4), input_y + Inches(0.1),
             Inches(2), Inches(0.3),
             font=BODY_FONT, size=9, italic=True, color=CHARCOAL_60)


def slide_07_crm_scoring(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "CRM intelligent", Inches(0.6), Inches(0.6))
    add_h2(slide,
           "Chaque prospect, une temperature. Tu sais qui rappeler en premier.",
           Inches(0.6), Inches(1.0), Inches(12), size=24)
    add_body(slide,
             "5 niveaux calcules automatiquement depuis la conversation SMS, le delai de reponse, et le budget.",
             Inches(0.6), Inches(2.0), Inches(12), Inches(0.6),
             size=14, color=CHARCOAL_60)

    # Card centre
    cw = Inches(8.5)
    ch = Inches(3.4)
    cx = Inches((SLIDE_W_IN - cw.inches) / 2)
    cy = Inches(2.85)
    add_card(slide, cx, cy, cw, ch)

    # Header card: nom + heat dots
    add_text(slide, "Marie-Eve Tremblay",
             cx + Inches(0.4), cy + Inches(0.35),
             Inches(5), Pt(28),
             font=BODY_FONT, size=18, bold=True, color=CHARCOAL)
    # heat dots a droite (4 sur 5)
    heat_w = add_heat_dots(slide, cx + cw - Inches(2.0), cy + Inches(0.4), level=4)

    # Mono caption
    add_text(slide, "R E C U   ·   V E R D U N   ·   4 2 5 K   ·   D U P L E X",
             cx + Inches(0.4), cy + Inches(0.95),
             Inches(7), Inches(0.3),
             font=MONO_FONT, size=8, bold=True, color=CHARCOAL_60)

    # Separateur
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 cx + Inches(0.4), cy + Inches(1.4),
                                 cw - Inches(0.8), Pt(1))
    sep.fill.solid()
    sep.fill.fore_color.rgb = BORDER
    sep.line.fill.background()

    # 3 mini-tiles
    tiles = [("Budget", "425 000 $"), ("Secteur", "Verdun, St-Henri"), ("Delai", "< 60 jours")]
    tile_w = (cw - Inches(0.8)) / 3
    for i, (label, value) in enumerate(tiles):
        tx = cx + Inches(0.4) + i * tile_w
        ty = cy + Inches(1.65)
        add_text(slide, label.upper(), tx, ty,
                 tile_w, Pt(18),
                 font=MONO_FONT, size=9, bold=True, color=CHARCOAL_60)
        add_text(slide, value, tx, ty + Inches(0.3),
                 tile_w, Pt(28),
                 font=BODY_FONT, size=16, bold=True, color=CHARCOAL)

    # Boutons CTA
    btn_y = cy + ch - Inches(0.75)
    add_button(slide, cx + Inches(0.4), btn_y,
               Inches(2.0), Inches(0.45), "Appeler maintenant", primary=True)
    add_button(slide, cx + Inches(2.6), btn_y,
               Inches(2.0), Inches(0.45), "Voir conversation", primary=False)

    # Legende heat scale en bas
    leg_y = Inches(6.55)
    add_text(slide,
             "F R O I D    ·    T I E D E    ·    C H A U D    ·    T R E S   C H A U D    ·    B R U L A N T",
             Inches(2.5), leg_y, Inches(8), Pt(20),
             font=MONO_FONT, size=8, bold=True, color=CHARCOAL_60,
             align=PP_ALIGN.CENTER)
    # 5 dots colores
    dot_d = Inches(0.18)
    total_dots_w = 5 * dot_d.inches + 4 * 0.30
    start_dx = (SLIDE_W_IN - total_dots_w) / 2
    for i in range(5):
        cxi = Inches(start_dx + dot_d.inches / 2 + i * (dot_d.inches + 0.30))
        cyi = leg_y + Inches(0.35)
        add_circle(slide, cxi, cyi, dot_d, fill=HEAT[i])

    add_footer(slide, 7)
    set_notes(slide,
              "Chaque prospect a une temperature. 5 niveaux : froid bleu jusqu'a "
              "brulant rouge. Calcule automatiquement, pas une note saisie a la "
              "main. Tu sais qui rappeler en premier des que tu sors de visite.")
    return slide


def slide_08_briefing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "Briefing quotidien", Inches(0.6), Inches(0.6))
    add_h2(slide,
           "Chaque matin a 7h30, ton plan de la journee arrive.",
           Inches(0.6), Inches(1.0), Inches(12), size=26)
    add_body(slide,
             "4 KPIs · les leads chauds · les nouveaux · les relances dues. Email + in-app.",
             Inches(0.6), Inches(1.85), Inches(12), Inches(0.6),
             size=14, color=CHARCOAL_60)

    # Big card dashboard
    cw = Inches(11.5)
    ch = Inches(4.3)
    cx = Inches((SLIDE_W_IN - cw.inches) / 2)
    cy = Inches(2.55)
    add_card(slide, cx, cy, cw, ch)

    # Header card: mark Klaris + greeting + date
    add_logo_mark(slide, cx + Inches(0.30), cy + Inches(0.20), Inches(0.55))
    add_text(slide, "Bon matin Joanel — voici ta journee",
             cx + Inches(1.0), cy + Inches(0.25),
             Inches(8), Pt(28),
             font=BODY_FONT, size=18, bold=True, color=CHARCOAL)
    add_text(slide, "M A R D I   5   M A I   ·   7 : 3 0",
             cx + cw - Inches(3.5), cy + Inches(0.3),
             Inches(3.0), Pt(18),
             font=MONO_FONT, size=9, bold=True, color=CHARCOAL_60,
             align=PP_ALIGN.RIGHT)

    # 4 KPI tiles
    kpis = [("12", "Leads actifs"), ("3", "Nouveaux qualifies"),
            ("5", "Relances dues"), ("2", "RDV aujourd'hui")]
    tile_w = Inches(2.5)
    tile_h = Inches(1.0)
    tile_gap = Inches(0.16)
    tiles_total = 4 * tile_w + 3 * tile_gap
    tile_start_x = cx + (cw - tiles_total) / 2
    tile_y = cy + Inches(1.05)
    for i, (n, lab) in enumerate(kpis):
        tx = tile_start_x + i * (tile_w + tile_gap)
        tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      tx, tile_y, tile_w, tile_h)
        tile.adjustments[0] = 0.10
        tile.fill.solid()
        tile.fill.fore_color.rgb = CREAM
        tile.line.fill.background()
        add_text(slide, n, tx + Inches(0.25), tile_y + Inches(0.12),
                 tile_w - Inches(0.5), Pt(40),
                 font=BODY_FONT, size=28, bold=True, color=TERRACOTTA)
        add_text(slide, lab, tx + Inches(0.25), tile_y + Inches(0.65),
                 tile_w - Inches(0.5), Pt(18),
                 font=BODY_FONT, size=10, bold=True, color=CHARCOAL)

    # 2 colonnes en bas
    col_y = cy + Inches(2.35)
    col_h = Inches(1.85)
    col_gap = Inches(0.4)
    col_w = (cw - Inches(0.8) - col_gap) / 2

    # Colonne 1 — A rappeler en priorite
    cx1 = cx + Inches(0.4)
    add_text(slide, "A rappeler en priorite",
             cx1, col_y, col_w, Pt(22),
             font=BODY_FONT, size=13, bold=True, color=CHARCOAL)
    rows = [
        ("Marie-Eve T.", 4, "Verdun, 425K, prete a signer"),
        ("Patrick B.", 5, "Plateau, 600K, hesite entre 2 plex"),
        ("Sandra L.", 3, "Rosemont, 380K, premier achat"),
    ]
    for i, (name, lvl, ctx) in enumerate(rows):
        ry = col_y + Inches(0.4) + i * Inches(0.42)
        add_text(slide, name, cx1, ry,
                 Inches(1.5), Inches(0.3),
                 font=BODY_FONT, size=11, bold=True, color=CHARCOAL)
        add_heat_dots(slide, cx1 + Inches(1.55), ry + Inches(0.05),
                      level=lvl, dot_d=Pt(7), gap=Pt(3))
        add_text(slide, ctx, cx1 + Inches(2.5), ry,
                 col_w - Inches(2.5), Inches(0.3),
                 font=BODY_FONT, size=10, color=CHARCOAL_60)

    # Colonne 2 — Aujourd'hui
    cx2 = cx + Inches(0.4) + col_w + col_gap
    add_text(slide, "Aujourd'hui",
             cx2, col_y, col_w, Pt(22),
             font=BODY_FONT, size=13, bold=True, color=CHARCOAL)
    appts = [
        ("14:00", "Visite 1234 rue Wellington"),
        ("16:30", "Suivi appel — David M."),
    ]
    for i, (t, label) in enumerate(appts):
        ay = col_y + Inches(0.4) + i * Inches(0.42)
        add_text(slide, t, cx2, ay,
                 Inches(1.0), Inches(0.3),
                 font=MONO_FONT, size=11, bold=True, color=TERRACOTTA)
        add_text(slide, label, cx2 + Inches(1.0), ay,
                 col_w - Inches(1.0), Inches(0.3),
                 font=BODY_FONT, size=11, color=CHARCOAL)

    add_footer(slide, 8)
    set_notes(slide,
              "Chaque matin a 7h30, ton plan arrive par email et dans l'app. 4 "
              "KPIs en haut. En dessous, les leads a rappeler en priorite avec leur "
              "temperature, et les RDV de la journee. Tu ouvres ton telephone, tu "
              "sais ou aller.")
    return slide


def slide_09_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "Fonctionnalites", Inches(0.6), Inches(0.6))
    add_h2(slide, "Concue avec des courtiers, pour des courtiers.",
           Inches(0.6), Inches(1.0), Inches(12), size=26)

    feats = [
        ("S", "Chatbot SMS",      "5 questions, ton naturel, signature « son assistante »"),
        ("T", "CRM scoring",      "5 niveaux de temperature, calcule auto"),
        ("R", "Relances auto",    "J+2 / J+5 / J+10, encadrees Loi 25 et CASL"),
        ("B", "Briefing 7h30",    "Email + in-app, 4 KPIs et plan du jour"),
        ("C", "Calendrier RDV",   "Strip 14 jours, sync iOS native"),
        ("V", "Voicemail intake", "Appel manque -> fiche prospect auto"),
        ("M", "Templates SMS",    "Placeholders {nom} {secteur} {budget}"),
        ("P", "Stats + PDF",      "Auto-genere chaque mois"),
    ]

    cell_w = Inches(5.8)
    cell_h = Inches(1.0)
    gap_x = Inches(0.4)
    gap_y = Inches(0.18)
    cols = 2
    grid_w = Emu(cols * cell_w + (cols - 1) * gap_x)
    start_x = Inches((SLIDE_W_IN - grid_w.inches) / 2)
    start_y = Inches(2.05)

    for i, (sym, title, body) in enumerate(feats):
        col = i % cols
        row = i // cols
        x = start_x + col * (cell_w + gap_x)
        y = start_y + row * (cell_h + gap_y)
        add_card(slide, x, y, cell_w, cell_h, fill=WHITE)
        # icone cercle terracotta-soft + lettre
        add_circle(slide, x + Inches(0.45), y + cell_h / 2,
                   Inches(0.45), fill=TERRACOTTA_SOFT)
        add_text(slide, sym, x + Inches(0.25), y + cell_h / 2 - Inches(0.15),
                 Inches(0.4), Inches(0.3),
                 font=BODY_FONT, size=14, bold=True, color=TERRACOTTA,
                 align=PP_ALIGN.CENTER)
        # titre + body
        add_text(slide, title, x + Inches(0.95), y + Inches(0.18),
                 cell_w - Inches(1.1), Pt(20),
                 font=BODY_FONT, size=14, bold=True, color=CHARCOAL)
        add_text(slide, body, x + Inches(0.95), y + Inches(0.55),
                 cell_w - Inches(1.1), Pt(22),
                 font=BODY_FONT, size=11, color=CHARCOAL_60, line_spacing=1.3)

    add_footer(slide, 9)
    set_notes(slide,
              "Voici l'inventaire des fonctionnalites livrees. Les 4 premieres "
              "sont les coeurs metier. Les 4 suivantes sont des outils du "
              "quotidien — calendrier, voicemail intake, templates SMS, stats et "
              "rapports PDF mensuels.")
    return slide


def slide_10_voicemail_calendrier(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "Intake complet", Inches(0.6), Inches(0.6))
    add_h2(slide,
           "Un appel manque ne se perd plus. Un RDV se prend en deux taps.",
           Inches(0.6), Inches(1.0), Inches(12), size=24)

    # 2 colonnes mockups
    col_w = Inches(5.8)
    col_h = Inches(4.3)
    gap = Inches(0.4)
    total = Emu(2 * col_w + gap)
    start_x = Inches((SLIDE_W_IN - total.inches) / 2)
    y = Inches(2.3)

    # GAUCHE — Voicemail
    lx = start_x
    add_card(slide, lx, y, col_w, col_h)
    add_text(slide, "A P P E L   M A N Q U E   ·   1 4 : 2 3   ·   3 8 s",
             lx + Inches(0.35), y + Inches(0.3),
             col_w - Inches(0.7), Pt(18),
             font=MONO_FONT, size=9, bold=True, color=TERRACOTTA)
    # Onde audio (10 mini-rect verticaux)
    wave_x = lx + Inches(0.35)
    wave_y = y + Inches(0.95)
    wave_h_max = Inches(0.55)
    bar_w = Inches(0.18)
    bar_gap = Inches(0.08)
    bar_heights = [0.30, 0.55, 0.42, 0.78, 0.95, 0.62, 0.88, 0.40, 0.70, 0.50,
                   0.65, 0.35, 0.82, 0.55, 0.45]
    for i, hf in enumerate(bar_heights):
        bx = wave_x + i * (bar_w + bar_gap)
        bh = Inches(wave_h_max.inches * hf)
        by = wave_y + (wave_h_max - bh) / 2
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     bx, by, bar_w, bh)
        bar.adjustments[0] = 0.5
        bar.fill.solid()
        bar.fill.fore_color.rgb = TERRACOTTA
        _set_shape_alpha(bar, 0.85)
        bar.line.fill.background()

    # Transcription italic
    transcript = ('« Allo, c\'est David Mercier, j\'ai vu votre annonce sur le '
                  'triplex Hochelaga, j\'aimerais avoir plus de details. Mon '
                  'numero... »')
    add_text(slide, transcript,
             lx + Inches(0.35), y + Inches(1.85),
             col_w - Inches(0.7), Inches(1.2),
             font=BODY_FONT, size=11, italic=True, color=CHARCOAL,
             line_spacing=1.4)

    # Mini-card fiche generee
    fc_y = y + col_h - Inches(0.95)
    fc_h = Inches(0.7)
    fc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                lx + Inches(0.35), fc_y,
                                col_w - Inches(0.7), fc_h)
    fc.adjustments[0] = 0.10
    fc.fill.solid()
    fc.fill.fore_color.rgb = CREAM
    fc.line.color.rgb = BORDER
    fc.line.width = Pt(0.75)
    add_text(slide, "FICHE GENEREE  ->  David Mercier · Hochelaga · triplex",
             lx + Inches(0.55), fc_y + Inches(0.13),
             col_w - Inches(2.0), Inches(0.25),
             font=MONO_FONT, size=8, bold=True, color=CHARCOAL)
    add_heat_dots(slide, lx + col_w - Inches(1.5), fc_y + Inches(0.40),
                  level=3, dot_d=Pt(7), gap=Pt(3))

    # DROITE — Calendrier
    rx = start_x + col_w + gap
    add_card(slide, rx, y, col_w, col_h)
    add_text(slide, "Mai 2026", rx + Inches(0.35), y + Inches(0.25),
             col_w - Inches(0.7), Pt(28),
             font=BODY_FONT, size=16, bold=True, color=CHARCOAL)

    # Strip 14 jours
    strip_x = rx + Inches(0.3)
    strip_y = y + Inches(0.95)
    cell_size = Inches(0.36)
    cell_gap = Inches(0.04)
    days = list(range(1, 15))
    selected_day = 7
    today_day = 5
    for i, d in enumerate(days):
        cxi = strip_x + i * (cell_size + cell_gap)
        cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      cxi, strip_y, cell_size, cell_size)
        cell.adjustments[0] = 0.20
        cell.fill.solid()
        if d == selected_day:
            cell.fill.fore_color.rgb = TERRACOTTA
            cell.line.fill.background()
            txt_color = WHITE
        elif d == today_day:
            cell.fill.fore_color.rgb = WHITE
            cell.line.color.rgb = TERRACOTTA
            cell.line.width = Pt(1.5)
            txt_color = TERRACOTTA
        else:
            cell.fill.fore_color.rgb = CREAM
            cell.line.color.rgb = BORDER
            cell.line.width = Pt(0.5)
            txt_color = CHARCOAL
        add_text(slide, str(d),
                 cxi, strip_y + Inches(0.07),
                 cell_size, Inches(0.2),
                 font=MONO_FONT, size=9, bold=True, color=txt_color,
                 align=PP_ALIGN.CENTER)
        # mini-dots RDV (variation)
        nb_rdv = [0, 1, 0, 2, 1, 0, 3, 1, 0, 0, 1, 2, 0, 1][i]
        for j in range(nb_rdv):
            dot_d2 = Pt(2.5)
            dx = cxi + cell_size / 2 - Pt(4) + j * Pt(3)
            dy = strip_y + cell_size - Pt(6)
            d2 = add_circle(slide, dx, dy, dot_d2,
                            fill=WHITE if d == selected_day else TERRACOTTA)

    # Liste 3 RDV jour selected
    list_y = strip_y + cell_size + Inches(0.5)
    add_text(slide, "Jeudi 7 mai",
             rx + Inches(0.35), list_y,
             col_w - Inches(0.7), Pt(20),
             font=BODY_FONT, size=12, bold=True, color=CHARCOAL_60)
    rdvs = [
        ("09:00", "Visite Verdun (Marie-Eve T.)"),
        ("14:00", "Open House Plateau"),
        ("16:30", "Suivi David M."),
    ]
    for i, (t, lab) in enumerate(rdvs):
        ay = list_y + Inches(0.45) + i * Inches(0.35)
        add_text(slide, t, rx + Inches(0.35), ay,
                 Inches(0.9), Inches(0.3),
                 font=MONO_FONT, size=10, bold=True, color=TERRACOTTA)
        add_text(slide, lab, rx + Inches(1.3), ay,
                 col_w - Inches(1.7), Inches(0.3),
                 font=BODY_FONT, size=11, color=CHARCOAL)

    # Caption mono bas
    add_text(slide,
             "S Y N C   I C L O U D   N A T I V E   ·   Z E R O   C O N F I G",
             rx + Inches(0.35), y + col_h - Inches(0.45),
             col_w - Inches(0.7), Inches(0.3),
             font=MONO_FONT, size=8, bold=True, color=CHARCOAL_60)

    add_footer(slide, 10)
    set_notes(slide,
              "Deux modules supplementaires. A gauche, le voicemail intake : "
              "quand un prospect appelle et que tu ne reponds pas, on transcrit en "
              "francais via Whisper, on qualifie via Claude, et la fiche est creee. "
              "A droite, le calendrier des RDV avec sync iOS native. Tout reste "
              "dans ton ecosysteme habituel.")
    return slide


def slide_11_stats_conformite(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "Stats et conformite", Inches(0.6), Inches(0.6))
    add_h2(slide,
           "Tu vois ta progression. La regie ne t'inquiete plus.",
           Inches(0.6), Inches(1.0), Inches(12), size=26)

    # Bloc gauche — chart
    lx = Inches(0.6)
    ly = Inches(2.2)
    lw = Inches(7.6)
    lh = Inches(4.6)
    add_card(slide, lx, ly, lw, lh)

    add_text(slide, "Performance — 6 derniers mois",
             lx + Inches(0.35), ly + Inches(0.25),
             lw - Inches(0.7), Pt(28),
             font=BODY_FONT, size=14, bold=True, color=CHARCOAL)

    # 4 mini KPI top row
    kpis = [("47", "Nouveaux"), ("28", "Qualifies"),
            ("8", "Conclus"), ("34%", "Taux close")]
    kpi_w = (lw - Inches(0.7)) / 4
    for i, (n, lab) in enumerate(kpis):
        kx = lx + Inches(0.35) + i * kpi_w
        ky = ly + Inches(0.85)
        add_text(slide, n, kx, ky,
                 kpi_w, Pt(28),
                 font=BODY_FONT, size=20, bold=True, color=TERRACOTTA)
        add_text(slide, lab, kx, ky + Inches(0.4),
                 kpi_w, Pt(18),
                 font=MONO_FONT, size=8, bold=True, color=CHARCOAL_60)

    # Bar chart 6 mois
    chart_x = lx + Inches(0.6)
    chart_y = ly + Inches(1.85)
    chart_w = lw - Inches(1.0)
    chart_h = Inches(1.6)

    months = ["DEC", "JAN", "FEV", "MAR", "AVR", "MAI"]
    # data: (nouveaux, qualifies, conclus) per mois
    data = [
        (28, 12, 3),
        (32, 15, 4),
        (38, 20, 5),
        (42, 24, 6),
        (45, 26, 7),
        (47, 28, 8),
    ]
    max_total = max(sum(d) for d in data)
    bar_w = chart_w / 9   # 6 bars + spacing
    bar_gap = bar_w * 0.5

    # Axe Y baseline
    base_y = chart_y + chart_h
    axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  chart_x, base_y, chart_w, Pt(0.75))
    axis.fill.solid()
    axis.fill.fore_color.rgb = BORDER
    axis.line.fill.background()

    for i, (n, q, c) in enumerate(data):
        bx = chart_x + i * (bar_w + bar_gap) + bar_gap
        # Total height
        total = n + q + c
        h_total = Inches(chart_h.inches * (total / max_total))
        by = base_y - h_total
        # nouveaux (cream)
        h_n = Inches(chart_h.inches * (n / max_total))
        b1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, bar_w, h_n)
        b1.fill.solid()
        b1.fill.fore_color.rgb = CREAM_DEEP
        b1.line.fill.background()
        # qualifies (terracotta soft)
        by2 = by + h_n
        h_q = Inches(chart_h.inches * (q / max_total))
        b2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by2, bar_w, h_q)
        b2.fill.solid()
        b2.fill.fore_color.rgb = TERRACOTTA
        _set_shape_alpha(b2, 0.45)
        b2.line.fill.background()
        # conclus (terracotta full)
        by3 = by2 + h_q
        h_c = Inches(chart_h.inches * (c / max_total))
        b3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by3, bar_w, h_c)
        b3.fill.solid()
        b3.fill.fore_color.rgb = TERRACOTTA
        b3.line.fill.background()
        # label mois
        add_text(slide, months[i],
                 bx - bar_gap / 2, base_y + Inches(0.06),
                 bar_w + bar_gap, Inches(0.25),
                 font=MONO_FONT, size=8, bold=True, color=CHARCOAL_60,
                 align=PP_ALIGN.CENTER)

    # Legende dots (sous month labels, au-dessus de la caption)
    leg_y = ly + lh - Inches(0.85)
    leg_items = [(CREAM_DEEP, "Nouveaux"), (TERRACOTTA_SOFT, "Qualifies"),
                 (TERRACOTTA, "Conclus")]
    leg_x = lx + Inches(0.6)
    for col, lab in leg_items:
        add_circle(slide, leg_x + Inches(0.1), leg_y + Inches(0.08),
                   Inches(0.12), fill=col)
        add_text(slide, lab, leg_x + Inches(0.25), leg_y,
                 Inches(1.5), Inches(0.25),
                 font=BODY_FONT, size=10, color=CHARCOAL_60)
        leg_x += Inches(1.6)

    # Bloc droite — Conformite
    rx = Inches(8.6)
    ry = Inches(2.2)
    rw = Inches(4.2)
    rh = Inches(4.6)
    add_card(slide, rx, ry, rw, rh)

    add_text(slide, "C O N F O R M I T E   B A K E E",
             rx + Inches(0.35), ry + Inches(0.3),
             rw - Inches(0.7), Inches(0.3),
             font=MONO_FONT, size=9, bold=True, color=TERRACOTTA)

    items = [
        ("OACIQ", "Regles metier integrees, pas un patch"),
        ("Loi 25", "Consentement, audit log, droits d'acces"),
        ("CASL", "Opt-out auto, horaires QC, frequence encadree"),
    ]
    iy = ry + Inches(1.0)
    for title, body in items:
        # Check vert
        chk = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     rx + Inches(0.35), iy,
                                     Inches(0.32), Inches(0.32))
        chk.fill.solid()
        chk.fill.fore_color.rgb = SUCCESS
        chk.line.fill.background()
        add_text(slide, "v",
                 rx + Inches(0.35), iy + Inches(0.02),
                 Inches(0.32), Inches(0.3),
                 font=BODY_FONT, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(slide, title,
                 rx + Inches(0.85), iy - Inches(0.02),
                 rw - Inches(1.2), Pt(22),
                 font=BODY_FONT, size=14, bold=True, color=CHARCOAL)
        add_text(slide, body,
                 rx + Inches(0.85), iy + Inches(0.30),
                 rw - Inches(1.2), Inches(0.7),
                 font=BODY_FONT, size=10, color=CHARCOAL_60, line_spacing=1.3)
        iy += Inches(1.05)

    # Mention legale bas
    add_text(slide,
             "DONNEES HEBERGEES CA-CENTRAL-1  ·  RLS POSTGRES",
             rx + Inches(0.35), ry + rh - Inches(0.45),
             rw - Inches(0.7), Inches(0.3),
             font=MONO_FONT, size=7, bold=True, color=CHARCOAL_60)

    # Caption sous la legende, au bas de la card
    add_text(slide,
             "R A P P O R T   P D F   A U T O   ·   C H A Q U E   1 E R   D U   M O I S",
             lx + Inches(0.35), ly + lh - Inches(0.4),
             lw - Inches(0.7), Inches(0.25),
             font=MONO_FONT, size=8, bold=True, color=CHARCOAL_60,
             align=PP_ALIGN.CENTER)

    add_footer(slide, 11)
    set_notes(slide,
              "Tu vois ta progression mois apres mois. 4 KPIs en haut. Chart 6 mois "
              "avec 3 niveaux : nouveaux, qualifies, conclus. Et un rapport PDF "
              "auto-genere chaque 1er du mois. A droite : la conformite. OACIQ, "
              "Loi 25, CASL — bakee dans l'archi, pas patchee apres.")
    return slide


def slide_12_pourquoi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "Momentum", Inches(0.6), Inches(0.6))
    add_h2(slide,
           "L'IA touche enfin les metiers terrain. Klaris est construit ici, par ici, pour ici.",
           Inches(0.6), Inches(1.0), Inches(12), size=22)

    # 2 colonnes
    col_w = Inches(5.7)
    gap = Inches(0.5)
    col_h = Inches(4.0)
    total = Emu(2 * col_w + gap)
    start_x = Inches((SLIDE_W_IN - total.inches) / 2)
    y = Inches(2.4)

    # Ligne verticale separation
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 start_x + col_w + gap / 2 - Pt(0.5), y,
                                 Pt(1), col_h)
    sep.fill.solid()
    sep.fill.fore_color.rgb = BORDER
    sep.line.fill.background()

    # Colonne G — Pourquoi maintenant
    lx = start_x
    add_text(slide, "Pourquoi maintenant",
             lx, y, col_w, Pt(28),
             font=BODY_FONT, size=18, bold=True, color=TERRACOTTA)
    pm = [
        ("Cout d'usage", "divise par 100 en 2 ans — l'IA est devenue accessible sur abonnement"),
        ("Adoption SMS", "95% des QC textent quotidiennement, 0% lit ses emails de courtiers a temps"),
        ("Pression reglementaire", "Loi 25, CASL, OACIQ — les outils non-conformes vont etre bannis"),
    ]
    iy = y + Inches(0.7)
    for title, body in pm:
        # Bullet dot terracotta
        add_circle(slide, lx + Inches(0.12), iy + Inches(0.18),
                   Inches(0.12), fill=TERRACOTTA)
        add_text(slide, title, lx + Inches(0.4), iy,
                 col_w - Inches(0.4), Pt(22),
                 font=BODY_FONT, size=14, bold=True, color=CHARCOAL)
        add_text(slide, body, lx + Inches(0.4), iy + Inches(0.32),
                 col_w - Inches(0.4), Inches(0.8),
                 font=BODY_FONT, size=12, color=CHARCOAL_60, line_spacing=1.35)
        iy += Inches(1.05)

    # Colonne D — Pourquoi nous
    rx = start_x + col_w + gap
    add_text(slide, "Pourquoi nous",
             rx, y, col_w, Pt(28),
             font=BODY_FONT, size=18, bold=True, color=TERRACOTTA)
    pn = [
        ("Methodologie JTBD", "4 personas valides, 12 patterns terrain, 0 hypothese"),
        ("Conformite d'abord", "bakee dans l'archi, pas patchee apres"),
        ("Klaris parle ton vocabulaire", "Centris, plex, hypotheque, secteur. Pas un CRM generique traduit."),
    ]
    iy = y + Inches(0.7)
    for title, body in pn:
        add_circle(slide, rx + Inches(0.12), iy + Inches(0.18),
                   Inches(0.12), fill=TERRACOTTA)
        add_text(slide, title, rx + Inches(0.4), iy,
                 col_w - Inches(0.4), Pt(22),
                 font=BODY_FONT, size=14, bold=True, color=CHARCOAL)
        add_text(slide, body, rx + Inches(0.4), iy + Inches(0.32),
                 col_w - Inches(0.4), Inches(0.8),
                 font=BODY_FONT, size=12, color=CHARCOAL_60, line_spacing=1.35)
        iy += Inches(1.05)

    add_footer(slide, 12)
    set_notes(slide,
              "Pourquoi maintenant : l'IA est accessible, le SMS est le canal n1 "
              "au QC, et la pression reglementaire bannit bientot les outils non "
              "conformes. Pourquoi nous : methodologie JTBD basee sur 4 courtiers "
              "interviewes en profondeur, conformite first, et vocabulaire local.")
    return slide


def slide_13_tarifs(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "Tarifs", Inches(0.6), Inches(0.6))
    add_h2(slide, "Un abonnement. Sans surprise.",
           Inches(0.6), Inches(1.0), Inches(12), size=28)

    # 2 cards
    card_w = Inches(4.6)
    card_h = Inches(4.7)
    gap = Inches(0.5)
    total = Emu(2 * card_w + gap)
    start_x = Inches((SLIDE_W_IN - total.inches) / 2)
    y = Inches(2.05)

    # Card SOLO (featured)
    sx = start_x
    # Card body
    add_card(slide, sx, y, card_w, card_h, fill=WHITE,
             border=TERRACOTTA, radius=0.06)
    # Bordure renforcee terracotta (rectangle externe pour effet 3pt)
    border = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    sx - Pt(2), y - Pt(2),
                                    card_w + Pt(4), card_h + Pt(4))
    border.adjustments[0] = 0.06
    border.fill.background()
    border.line.color.rgb = TERRACOTTA
    border.line.width = Pt(2.5)

    # Badge "RECOMMANDE POUR TOI" — style site (top: -12px, right: 24px)
    badge_w = Inches(2.0)
    badge_h = Inches(0.32)
    bx = sx + card_w - badge_w - Inches(0.25)
    by = y - badge_h / 2 - Pt(2)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   bx, by, badge_w, badge_h)
    badge.adjustments[0] = 0.5
    badge.fill.solid()
    badge.fill.fore_color.rgb = TERRACOTTA
    badge.line.fill.background()
    add_text(slide, "RECOMMANDE POUR TOI",
             bx, by + Inches(0.06),
             badge_w, Inches(0.3),
             font=MONO_FONT, size=8, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    # Titre Solo
    add_text(slide, "Solo", sx + Inches(0.4), y + Inches(0.5),
             card_w - Inches(0.8), Pt(36),
             font=BODY_FONT, size=22, bold=True, color=CHARCOAL)
    # Prix
    add_text(slide, "100 $ CAD",
             sx + Inches(0.4), y + Inches(1.05),
             card_w - Inches(0.8), Pt(50),
             font=BODY_FONT, size=32, bold=True, color=TERRACOTTA)
    add_text(slide, "/ mois",
             sx + Inches(0.4), y + Inches(1.85),
             card_w - Inches(0.8), Pt(20),
             font=BODY_FONT, size=12, color=CHARCOAL_60)

    # Features
    feats = ["Chatbot SMS illimite",
             "CRM avec scoring temperature",
             "Relances automatiques J+2 / J+5 / J+10",
             "Briefing quotidien 7h30",
             "Bilingue FR / EN"]
    fy = y + Inches(2.25)
    for f in feats:
        add_text(slide, "v", sx + Inches(0.4), fy,
                 Inches(0.25), Inches(0.28),
                 font=BODY_FONT, size=12, bold=True, color=TERRACOTTA)
        add_text(slide, f, sx + Inches(0.7), fy,
                 card_w - Inches(1.0), Inches(0.28),
                 font=BODY_FONT, size=11, color=CHARCOAL)
        fy += Inches(0.30)

    # CTA primary
    add_button(slide, sx + Inches(0.5), y + card_h - Inches(0.65),
               card_w - Inches(1.0), Inches(0.45),
               "Commencer l'essai", primary=True)

    # Card AGENCE (discrete)
    ax = start_x + card_w + gap
    add_card(slide, ax, y, card_w, card_h, fill=CREAM,
             border=BORDER, radius=0.06, shadow=False)

    add_text(slide, "Agence", ax + Inches(0.4), y + Inches(0.5),
             card_w - Inches(0.8), Pt(36),
             font=BODY_FONT, size=22, bold=True, color=CHARCOAL)
    add_text(slide, "200 $ CAD",
             ax + Inches(0.4), y + Inches(1.05),
             card_w - Inches(0.8), Pt(50),
             font=BODY_FONT, size=32, bold=True, color=CHARCOAL)
    add_text(slide, "/ courtier / mois",
             ax + Inches(0.4), y + Inches(1.85),
             card_w - Inches(0.8), Pt(20),
             font=BODY_FONT, size=12, color=CHARCOAL_60)

    feats_a = ["Tout du Solo",
               "Dashboard direction",
               "Templates personnalises",
               "SLA 99.9% + onboarding dedie"]
    fy = y + Inches(2.25)
    for f in feats_a:
        add_text(slide, "v", ax + Inches(0.4), fy,
                 Inches(0.25), Inches(0.28),
                 font=BODY_FONT, size=12, bold=True, color=CHARCOAL_60)
        add_text(slide, f, ax + Inches(0.7), fy,
                 card_w - Inches(1.0), Inches(0.28),
                 font=BODY_FONT, size=11, color=CHARCOAL)
        fy += Inches(0.30)

    add_button(slide, ax + Inches(0.5), y + card_h - Inches(0.65),
               card_w - Inches(1.0), Inches(0.45),
               "Nous parler", primary=False)

    # Caption sous les cards (au-dessus du footer)
    add_text(slide,
             "P A S   D ' E N G A G E M E N T   ·   A N N U L A T I O N   E N   U N   C L I C   ·   E S S A I   1 4   J O U R S",
             Inches(0.6), Inches(6.9), Inches(SLIDE_W_IN - 1.2), Inches(0.25),
             font=MONO_FONT, size=8, bold=True, color=CHARCOAL_60,
             align=PP_ALIGN.CENTER)

    add_footer(slide, 13)
    set_notes(slide,
              "Solo : 100$ par mois. Tout inclus. Pas de frais de mise en route, "
              "annulable a tout moment, hebergement Canada inclus. Agence : 200$ "
              "par courtier par mois pour les franchises. Compare a 2500-3000$ "
              "pour une assistante temps plein QC, c'est divise par 25-30.")
    return slide


def slide_14_temoignage_contact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # === TOP — Temoignage ===
    add_eyebrow(slide, "Pilote terrain", Inches(0.6), Inches(0.6))

    # Gros guillemet decoratif
    add_text(slide, "«",
             Inches(0.6), Inches(0.9),
             Inches(1.5), Inches(1.5),
             font=BODY_FONT, size=110, bold=True, color=TERRACOTTA_SOFT,
             line_spacing=1.0)

    quote = ("Une IA qui genere une offre d'achat en un rien de temps. Pendant "
             "que je suis en visite, Klaris parle a mes prospects et me sort une "
             "fiche prete. Avant, je rentrais au bureau et je perdais 30-60 "
             "minutes par offre. La, c'est secondes.")
    add_text(slide, quote,
             Inches(2.0), Inches(1.3),
             Inches(SLIDE_W_IN - 2.6), Inches(2.3),
             font=BODY_FONT, size=18, italic=True, color=CHARCOAL,
             line_spacing=1.45)

    # Attribution (sans letterspaced pour eviter wrap)
    add_text(slide,
             "JOANEL D.  ·  COURTIER IMMOBILIER  ·  GRAND MONTREAL  ·  90% CLOSE RATE  ·  PILOTE KLARIS",
             Inches(2.0), Inches(3.65),
             Inches(SLIDE_W_IN - 2.6), Inches(0.4),
             font=MONO_FONT, size=9, bold=True, color=TERRACOTTA)

    # Separateur horizontal
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), Inches(4.3),
                                 Inches(SLIDE_W_IN - 1.2), Pt(0.75))
    sep.fill.solid()
    sep.fill.fore_color.rgb = BORDER
    sep.line.fill.background()

    # === BOTTOM — Contact ===
    add_eyebrow(slide, "Prochaines etapes", Inches(0.6), Inches(4.55))
    add_h2(slide, "Reservons 15 minutes.",
           Inches(0.6), Inches(4.95), Inches(12), size=26)
    add_text(slide,
             "Demo personnalisee, sur ton cas reel. Pas de pression. Pas de carte de credit.",
             Inches(0.6), Inches(5.6), Inches(SLIDE_W_IN - 1.2), Inches(0.5),
             font=BODY_FONT, size=14, color=CHARCOAL_60)

    # Bloc CTA terracotta
    cta_w = Inches(7.0)
    cta_h = Inches(0.9)
    cx = Inches((SLIDE_W_IN - cta_w.inches) / 2)
    cy = Inches(6.10)
    cta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 cx, cy, cta_w, cta_h)
    cta.adjustments[0] = 0.18
    cta.fill.solid()
    cta.fill.fore_color.rgb = TERRACOTTA
    cta.line.fill.background()
    add_text(slide, "contact@klarisapp.ai",
             cx, cy + Inches(0.15),
             cta_w, Pt(28),
             font=BODY_FONT, size=17, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, "klarisapp.ai",
             cx, cy + Inches(0.52),
             cta_w, Pt(22),
             font=MONO_FONT, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    # Footer custom (sous le CTA, marge claire)
    add_text(slide, "Klaris — une marque Next Move · 2026",
             Inches(0.6), Inches(SLIDE_H_IN - 0.32),
             Inches(SLIDE_W_IN - 1.2), Inches(0.25),
             font=MONO_FONT, size=8, bold=True, color=CHARCOAL_60,
             align=PP_ALIGN.CENTER)

    set_notes(slide,
              "Temoignage de Joanel, pilote actuel. Lis lentement, laisse "
              "respirer. Puis : on ne te demande pas de t'engager. On te demande "
              "15 minutes. Demo personnalisee sur ton cas reel. Pas de pression, "
              "pas de carte de credit. Mon email est sur l'ecran. A toi de jouer.")
    return slide


# ----------------------------------------------------------------------------
# Build
# ----------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slide_01_cover(prs)
    slide_02_probleme(prs)
    slide_03_cout_inaction(prs)
    slide_04_solution(prs)
    slide_05_comment(prs)
    slide_06_demo_sms(prs)
    slide_07_crm_scoring(prs)
    slide_08_briefing(prs)
    slide_09_features(prs)
    slide_10_voicemail_calendrier(prs)
    slide_11_stats_conformite(prs)
    slide_12_pourquoi(prs)
    slide_13_tarifs(prs)
    slide_14_temoignage_contact(prs)

    out = Path(__file__).parent / "Klaris_AvantVente_Solo_FR.pptx"
    prs.save(out)
    print(f"OK -> {out}")


if __name__ == "__main__":
    build()
