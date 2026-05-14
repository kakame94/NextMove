"""Klaris — Slides annexes partagees (acronymes, glossaire, sources).

Helpers reutilisables par build_klaris_investor_deck.py et
build_klaris_gtm_deck.py. Brand tokens importes de build_klaris_deck.py.
"""

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from build_klaris_deck import (
    TERRACOTTA, TERRACOTTA_STRONG, TERRACOTTA_SOFT,
    CREAM, CHARCOAL, CHARCOAL_60, WHITE, BORDER,
    BODY_FONT, MONO_FONT, SLIDE_W_IN, SLIDE_H_IN,
    add_bg, add_text, add_card,
    add_footer, set_notes,
)
from build_klaris_investor_deck import add_section_header


# Acronymes regroupes par categorie. ASCII-safe (sans accents) pour eviter
# tout probleme de rendu dans PowerPoint / Keynote.

ACRONYMES_IMMOBILIER = [
    ("OACIQ", "Organisme d'autoreglementation du courtage immobilier QC"),
    ("APCIQ", "Association professionnelle des courtiers immobiliers QC"),
    ("ACI / CREA", "Association canadienne de l'immeuble"),
    ("RECO", "Real Estate Council of Ontario"),
    ("BCFSA", "BC Financial Services Authority"),
    ("FARCIQ", "Fonds d'assurance responsabilite courtage immo QC"),
    ("MLS", "Multiple Listing Service (Centris au QC)"),
    ("DDF", "Data Distribution Facility (flux MLS CREA)"),
    ("TRESA", "Trust in Real Estate Services Act (Ontario)"),
    ("TRREB", "Toronto Regional Real Estate Board"),
    ("GCI", "Gross Commission Income (revenu brut commissions)"),
    ("DOM", "Days on Market (delai de vente moyen)"),
    ("UFC", "Unites de formation continue (OACIQ)"),
    ("RMR", "Region metropolitaine de recensement"),
]

ACRONYMES_LEGAL_TECH = [
    ("Loi 25 / PL64", "Protection des renseignements personnels (QC)"),
    ("CAI", "Commission d'acces a l'information (QC)"),
    ("EFVP", "Evaluation des facteurs relatifs a la vie privee"),
    ("SAP", "Sanctions administratives pecuniaires (Loi 25)"),
    ("CASL / LCAP", "Canadian Anti-Spam Legislation"),
    ("CRTC", "Conseil radiodiffusion + telecoms canadiennes"),
    ("A2P", "Application-to-Person (SMS commerciaux)"),
    ("API", "Application Programming Interface"),
    ("LLM", "Large Language Model (Claude, GPT)"),
    ("RAG", "Retrieval-Augmented Generation"),
    ("RLS", "Row-Level Security (Postgres / Supabase)"),
    ("CRM", "Customer Relationship Management"),
    ("E&O", "Errors & Omissions (assurance prof.)"),
    ("ToS", "Terms of Service"),
    ("DPO", "Data Protection Officer"),
    ("SaaS", "Software as a Service"),
]

ACRONYMES_FINANCE_GTM = [
    ("TAM / SAM / SOM", "Total / Serviceable / Obtainable Addressable Market"),
    ("ARR", "Annual Recurring Revenue"),
    ("MRR", "Monthly Recurring Revenue"),
    ("ARPU", "Average Revenue Per User"),
    ("CAC", "Customer Acquisition Cost"),
    ("LTV", "Lifetime Value (valeur a vie)"),
    ("NRR", "Net Revenue Retention"),
    ("PMF", "Product-Market Fit"),
    ("ICP", "Ideal Customer Profile"),
    ("MQL", "Marketing Qualified Lead"),
    ("WTP", "Willingness To Pay"),
    ("GTM", "Go-To-Market"),
    ("CPM / CPC / CTR", "Cost per Mille / Click / Click-Through Rate"),
    ("B2B", "Business-to-Business"),
    ("SLA", "Service Level Agreement"),
    ("ROI", "Return on Investment"),
    ("VP", "Vice-President"),
    ("BDC / PME / PSCE", "Bailleurs / subventions QC (BDC, MEI, PME MTL)"),
]


def _draw_acronym_column(slide, x, y, w, h, title, items):
    """Trois cartes-colonnes : titre + table compact acronymes."""
    add_card(slide, x, y, w, h, fill=WHITE, border=BORDER, radius=0.10)
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x, y, w, Inches(0.42))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TERRACOTTA
    bar.line.fill.background()
    add_text(slide, title.upper(),
             x + Inches(0.20), y + Inches(0.08),
             w - Inches(0.4), Inches(0.35),
             font=MONO_FONT, size=9, bold=True, color=WHITE,
             line_spacing=1.0)

    # Items list
    row_h = Inches(0.27)
    start_y = y + Inches(0.55)
    for i, (acro, definition) in enumerate(items):
        ry = start_y + i * row_h
        # Acronym (left, terracotta, monospace bold)
        add_text(slide, acro,
                 x + Inches(0.18), ry,
                 Inches(1.55), row_h,
                 font=MONO_FONT, size=8, bold=True, color=TERRACOTTA_STRONG,
                 line_spacing=1.15)
        # Definition (right, charcoal, small)
        add_text(slide, definition,
                 x + Inches(1.75), ry,
                 w - Inches(1.95), row_h,
                 font=BODY_FONT, size=7.5, color=CHARCOAL_60,
                 line_spacing=1.15)


def slide_annexe_acronymes(prs, page_num, total_slides):
    """Slide annexe 3-colonnes : Immobilier · Legal/Tech · Finance/GTM."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Annexe",
                       "Glossaire — acronymes utilises dans le deck.")

    col_w = Inches(4.10)
    col_h = Inches(4.95)
    col_gap = Inches(0.10)
    start_x = Inches(0.7)
    y = Inches(2.0)

    columns = [
        ("Immobilier & marche", ACRONYMES_IMMOBILIER),
        ("Legal · Tech · Conformite", ACRONYMES_LEGAL_TECH),
        ("Finance · SaaS · GTM", ACRONYMES_FINANCE_GTM),
    ]
    for i, (title, items) in enumerate(columns):
        x = start_x + i * (col_w + col_gap)
        _draw_acronym_column(slide, x, y, col_w, col_h, title, items)

    # Footer note
    add_text(slide,
             "Source complete : Etude de marche v3 §9 Bibliographie · "
             "OACIQ.com · cai.gouv.qc.ca · crtc.gc.ca · CREA.ca",
             Inches(0.7), Inches(7.10),
             Inches(SLIDE_W_IN - 1.4), Inches(0.3),
             font=MONO_FONT, size=8, italic=True, color=CHARCOAL_60,
             line_spacing=1.3)

    add_footer(slide, page_num, total_slides)
    set_notes(slide,
              "Annexe glossaire. Repris dans tous les decks Klaris. "
              "Distribuer aussi en tant que page PDF separee dans data room.")
    return slide


# ============================================================================
# Slide annexe 2 — Metriques produit / SaaS expliquees
# ============================================================================
#
# Audience : equipe interne avec niveau heterogene sur metriques SaaS B2B.
# Format : 10 cartes 2x5 grid · acronyme + def plain FR + formule + cible Klaris.
# ============================================================================

METRIQUES = [
    {
        "acro": "ARR / MRR",
        "nom": "Annual / Monthly Recurring Revenue",
        "def": "Revenus recurrents previsibles d'un produit par abonnement.",
        "formule": "MRR = somme abonnements actifs · ARR = MRR × 12",
        "cible": "Klaris M12 : 12,5 K$ MRR · 150 K$ ARR",
    },
    {
        "acro": "ARPU",
        "nom": "Average Revenue Per User",
        "def": "Revenu moyen genere par client sur une periode (an).",
        "formule": "ARPU = ARR total / nb clients actifs",
        "cible": "Klaris mix 70/30 : 1 500 $/an (1 200 Solo + 2 400 Agence)",
    },
    {
        "acro": "CAC",
        "nom": "Customer Acquisition Cost",
        "def": "Cout total marketing + ventes pour acquerir un client.",
        "formule": "CAC = (depenses marketing + ventes) / nouveaux clients",
        "cible": "Klaris cible : < 1 200 $ blended (Meta + outbound)",
    },
    {
        "acro": "LTV",
        "nom": "Lifetime Value",
        "def": "Revenu total qu'un client genere avant churn (depart).",
        "formule": "LTV = ARPU × duree moyenne retention (annees)",
        "cible": "Klaris : 3 000 $ (1 500 $/an × 2 ans retention)",
    },
    {
        "acro": "LTV : CAC",
        "nom": "Ratio rentabilite acquisition",
        "def": "Combien rapporte un client compare a ce qu'il a coute.",
        "formule": "Ratio = LTV / CAC · standard SaaS sain = 3:1+",
        "cible": "Klaris : 2,5:1 (acceptable bootstrap · ameliorer via referral)",
    },
    {
        "acro": "Payback",
        "nom": "Payback Period",
        "def": "Nb de mois pour recuperer le CAC via abonnement client.",
        "formule": "Payback = CAC / MRR moyen par client",
        "cible": "Klaris : 9,6 mois (1 200 $ / 125 $/mo) · limite haute",
    },
    {
        "acro": "Churn logo",
        "nom": "Taux d'attrition clients",
        "def": "% clients qui annulent leur abonnement chaque mois.",
        "formule": "Churn = clients perdus / clients actifs debut periode",
        "cible": "Klaris : < 3 %/mois (PME SaaS B2B sain) · alerte > 5 %",
    },
    {
        "acro": "NRR",
        "nom": "Net Revenue Retention",
        "def": "Revenus conserves + expansions - downgrades - churns par cohorte.",
        "formule": "NRR = (MRR cohorte fin - downgrades + upsells) / MRR debut",
        "cible": "Klaris cible M24 : > 100 % (upsell Solo Pro → Agence)",
    },
    {
        "acro": "MQL / SQL",
        "nom": "Marketing / Sales Qualified Lead",
        "def": "Prospect qualifie par marketing (MQL) puis valide ventes (SQL).",
        "formule": "Conversion MQL → SQL → client = entonnoir B2B",
        "cible": "Klaris : 8-12 MQL/sem · 10 % conversion client",
    },
    {
        "acro": "PMF",
        "nom": "Product-Market Fit",
        "def": "Moment ou marche tire le produit (demande > capacite servir).",
        "formule": "Proxy = NRR > 100 % + churn < 3 % sur cohorte 3-6 mois",
        "cible": "Klaris jalon : valide M6 (50 clients · retention nette 90 %+)",
    },
]


def _draw_metric_card(slide, x, y, w, h, metric):
    """Carte metrique : acronyme + nom + def + formule + cible."""
    add_card(slide, x, y, w, h, fill=WHITE, border=BORDER, radius=0.10)

    # Acronyme + nom dans le header
    add_text(slide, metric["acro"],
             x + Inches(0.22), y + Inches(0.15),
             w - Inches(0.4), Inches(0.35),
             font=BODY_FONT, size=15, bold=True, color=TERRACOTTA_STRONG,
             line_spacing=1.1)
    add_text(slide, metric["nom"],
             x + Inches(0.22), y + Inches(0.45),
             w - Inches(0.4), Inches(0.25),
             font=MONO_FONT, size=8, color=CHARCOAL_60,
             line_spacing=1.15)

    # Def
    add_text(slide, metric["def"],
             x + Inches(0.22), y + Inches(0.75),
             w - Inches(0.4), Inches(0.5),
             font=BODY_FONT, size=8.5, color=CHARCOAL,
             line_spacing=1.30)

    # Formule (italic)
    add_text(slide, "Formule : " + metric["formule"],
             x + Inches(0.22), y + Inches(1.30),
             w - Inches(0.4), Inches(0.45),
             font=BODY_FONT, size=8, italic=True, color=CHARCOAL_60,
             line_spacing=1.30)

    # Cible Klaris (terracotta soft band)
    band_y = y + h - Inches(0.55)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, band_y, w, Inches(0.55))
    band.fill.solid()
    band.fill.fore_color.rgb = TERRACOTTA_SOFT
    band.line.fill.background()
    add_text(slide, "CIBLE KLARIS",
             x + Inches(0.22), band_y + Inches(0.06),
             w - Inches(0.4), Inches(0.2),
             font=MONO_FONT, size=7, bold=True, color=TERRACOTTA_STRONG,
             line_spacing=1.0)
    add_text(slide, metric["cible"],
             x + Inches(0.22), band_y + Inches(0.24),
             w - Inches(0.4), Inches(0.3),
             font=BODY_FONT, size=8.5, bold=True, color=TERRACOTTA_STRONG,
             line_spacing=1.20)


def slide_annexe_metriques(prs, page_num, total_slides):
    """Slide annexe : 10 metriques cle expliquees · 5 col × 2 lignes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Annexe · Metriques cle",
                       "Definitions + formules + cibles. Aligner l'equipe.")

    # 5 colonnes × 2 lignes
    cols = 5
    rows = 2
    card_w = Inches(2.42)
    card_h = Inches(2.30)
    gap_x = Inches(0.08)
    gap_y = Inches(0.12)
    start_x = Inches(0.7)
    start_y = Inches(2.05)

    for i, metric in enumerate(METRIQUES):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        _draw_metric_card(slide, x, y, card_w, card_h, metric)

    # Bottom note — link to data sources
    add_text(slide,
             "Ces metriques sont surveillees chaque lundi 10h en sync hebdo. "
             "Escalation : churn > 5 %/mois ou CAC > 1 500 $ pendant 2 semaines.",
             Inches(0.7), Inches(7.10),
             Inches(SLIDE_W_IN - 1.4), Inches(0.3),
             font=MONO_FONT, size=8, italic=True, color=CHARCOAL_60,
             line_spacing=1.3)

    add_footer(slide, page_num, total_slides)
    set_notes(slide,
              "Annexe pedagogique. Permet de niveler l'equipe sur metriques SaaS B2B. "
              "A relire chaque trimestre. Cibles Klaris extraites etude v3 §4 + §7.4.")
    return slide
