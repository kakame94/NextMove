"""Klaris — Deck investisseur / partenaire (FR). Generator python-pptx.

Output: Klaris_Investisseur_FR.pptx (13 slides, 16:9).
Source donnees: Etude de marche v3 + bibliographie Section 9.
Audience: investisseurs, partenaires strategiques (Centris, OACIQ, APCIQ),
banques (BDC, MEI Quebec, PME MTL), comites fondateurs.

Annexes (PDF separe, pas dans deck pitch live):
- Stress test churn 3/5/7% × ARR M24
- P&L Y1-Y3 detaille
- Cap table
- ToS / Limitation of Liability draft
- Etude de marche v3 complete

Helpers + brand tokens importes de build_klaris_deck.py (DRY, sync auto).
Aucune dependance PNG (mascotte = shape, logo = wordmark texte).
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from build_klaris_deck import (
    TERRACOTTA, TERRACOTTA_STRONG, TERRACOTTA_SOFT,
    CREAM, CREAM_DEEP, CHARCOAL, CHARCOAL_60, WHITE, BORDER,
    SUCCESS, WARN_AMBER, DESTRUCTIVE, INFO,
    BODY_FONT, MONO_FONT, SLIDE_W_IN, SLIDE_H_IN,
    add_bg, add_text, add_multitext, add_eyebrow,
    add_h1, add_h2, add_body, add_card,
    _set_shape_alpha,
    add_grid_pattern, add_circle, add_mascot,
    add_footer, set_notes,
)

TOTAL_SLIDES = 15


# ----------------------------------------------------------------------------
# Helpers specifiques deck investisseur
# ----------------------------------------------------------------------------

def add_section_header(slide, eyebrow, title, *, x=Inches(0.7), y=Inches(0.55)):
    add_eyebrow(slide, eyebrow, x, y)
    add_h1(slide, title, x, y + Inches(0.55), Inches(SLIDE_W_IN - 1.4), size=32)


def add_kpi_card(slide, x, y, w, h, label, value, sub, *,
                 value_color=TERRACOTTA, fill=WHITE):
    add_card(slide, x, y, w, h, fill=fill, border=BORDER, radius=0.10)
    add_text(slide, label.upper(), x + Inches(0.25), y + Inches(0.22),
             w - Inches(0.5), Pt(20),
             font=MONO_FONT, size=9, bold=True, color=CHARCOAL_60,
             line_spacing=1.0)
    add_text(slide, value, x + Inches(0.25), y + Inches(0.50),
             w - Inches(0.5), Pt(48),
             font=BODY_FONT, size=28, bold=True, color=value_color,
             line_spacing=1.0)
    add_text(slide, sub, x + Inches(0.25), y + h - Inches(0.55),
             w - Inches(0.5), Inches(0.4),
             font=BODY_FONT, size=10, color=CHARCOAL_60,
             line_spacing=1.25)


def add_table_row(slide, cells, x, y, col_widths, h, *,
                  header=False, bg=None):
    """cells: list[str]. col_widths: list[Inches]. Draws bordered row."""
    cur = x
    for i, txt in enumerate(cells):
        w = col_widths[i]
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cur, y, w, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg if bg else (CREAM_DEEP if header else WHITE)
        rect.line.color.rgb = BORDER
        rect.line.width = Pt(0.5)
        tf = rect.text_frame
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(4)
        tf.margin_bottom = Pt(4)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = txt
        f = run.font
        f.name = BODY_FONT
        f.size = Pt(10)
        f.bold = header
        f.color.rgb = CHARCOAL if header else CHARCOAL
        cur += w


def add_bullet_line(slide, dot_color, label, value, x, y, w):
    """Inline row: colored dot + label + value (right)."""
    add_circle(slide, x + Inches(0.08), y + Inches(0.12),
               Inches(0.16), fill=dot_color)
    add_text(slide, label, x + Inches(0.35), y, w - Inches(2.0), Inches(0.3),
             font=BODY_FONT, size=12, color=CHARCOAL, line_spacing=1.2)
    add_text(slide, value, x + w - Inches(2.0), y, Inches(2.0), Inches(0.3),
             font=MONO_FONT, size=12, bold=True, color=TERRACOTTA_STRONG,
             align=PP_ALIGN.RIGHT, line_spacing=1.2)


def add_wordmark(slide, x, y, *, size=14, color=CHARCOAL):
    add_text(slide, "K L A R I S", x, y, Inches(2), Pt(size * 1.6),
             font=MONO_FONT, size=size, bold=True, color=color,
             line_spacing=1.0)


# ----------------------------------------------------------------------------
# Slides
# ----------------------------------------------------------------------------

def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_grid_pattern(slide, Inches(0), Inches(0),
                     Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
                     step_in=0.45, alpha=0.30)

    # Glow behind mascot (right side)
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(8.4), Inches(1.5),
                                  Inches(4.5), Inches(4.5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = TERRACOTTA
    _set_shape_alpha(glow, 0.10)
    glow.line.fill.background()

    add_mascot(slide, Inches(10.65), Inches(3.75), Inches(2.6))

    # Wordmark + version
    add_wordmark(slide, Inches(0.7), Inches(0.7), size=16, color=CHARCOAL)
    add_text(slide, "DECK INVESTISSEUR · v1.0 · Q2 2026",
             Inches(0.7), Inches(1.15), Inches(6), Pt(18),
             font=MONO_FONT, size=10, color=CHARCOAL_60)

    # Hero title
    add_text(slide, "L'adjointe IA",
             Inches(0.7), Inches(2.6), Inches(8.2), Inches(1.2),
             font=BODY_FONT, size=56, bold=True, color=CHARCOAL,
             line_spacing=1.05)
    add_text(slide, "des courtiers immobiliers du Quebec.",
             Inches(0.7), Inches(3.55), Inches(8.2), Inches(1.0),
             font=BODY_FONT, size=42, bold=True, color=TERRACOTTA,
             line_spacing=1.05)

    # Sub
    add_text(slide,
             "Micro-SaaS B2B vertical. Bootstrap. Marche QC : SAM 17,85 M$ ARR. "
             "Cible 24 mois : 400-600 courtiers · 600-900 K$ ARR.",
             Inches(0.7), Inches(4.85), Inches(8.5), Inches(1.6),
             font=BODY_FONT, size=15, color=CHARCOAL_60, line_spacing=1.45)

    # Footer
    add_text(slide, "Dennis · Eliot · Walkens · Seydou — alanmanou.consulting@gmail.com",
             Inches(0.7), Inches(6.85), Inches(10), Pt(20),
             font=MONO_FONT, size=10, color=CHARCOAL_60)

    set_notes(slide,
              "Cover deck investisseur. Audience: BDC, MEI, PME MTL, partenaires "
              "strategiques (Centris, OACIQ, APCIQ), comite fondateurs. "
              "Position: Micro-SaaS rentable Bootstrap, pas pitch VC Serie A.")


def slide_02_probleme(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Le probleme",
                       "L'administratif tue les ventes. Les courtiers QC saturent.")

    # 3 KPI problem cards
    cards = [
        ("Delai de reponse au lead", "1 h+", "Prospect refroidi avant 1er contact.",
         DESTRUCTIVE),
        ("Adjointe humaine bilingue QC", "45-55 K$/an", "+ avantages = 4 500 $/mois charge.",
         WARN_AMBER),
        ("Adoption CRM franchise", "< 30 %", "Saisie redondante. CRM reste vide.",
         INFO),
    ]
    card_w = Inches(3.85)
    card_h = Inches(2.5)
    gap = Inches(0.20)
    start_x = Inches(0.7)
    y = Inches(2.4)
    for i, (lbl, val, sub, col) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_kpi_card(slide, x, y, card_w, card_h, lbl, val, sub, value_color=col)

    # Citation
    add_card(slide, Inches(0.7), Inches(5.3),
             Inches(SLIDE_W_IN - 1.4), Inches(1.5),
             fill=TERRACOTTA_SOFT, border=None, radius=0.12, shadow=False)
    add_text(slide,
             "« La difference entre moi et les courtiers a 40 M$, "
             "c'est juste la grosseur de l'equipe admin. »",
             Inches(1.0), Inches(5.55), Inches(SLIDE_W_IN - 2.0), Inches(0.8),
             font=BODY_FONT, size=16, italic=True, bold=True, color=TERRACOTTA_STRONG,
             line_spacing=1.3)
    add_text(slide, "— Joanel, courtier independant · beta-testeur Klaris",
             Inches(1.0), Inches(6.25), Inches(SLIDE_W_IN - 2.0), Inches(0.4),
             font=MONO_FONT, size=10, color=TERRACOTTA_STRONG, line_spacing=1.2)

    add_footer(slide, 2, TOTAL_SLIDES)
    set_notes(slide,
              "Problem framing en 3 chiffres durs. Source: salaire adjointe QC = "
              "Statistique Canada / Guichet Emplois. CRM adoption < 30% = NAR "
              "Technology Survey 2024.")


def slide_03_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "La solution",
                       "Klaris agit. Le courtier ferme.")

    # Killer message banner
    add_card(slide, Inches(0.7), Inches(2.0),
             Inches(SLIDE_W_IN - 1.4), Inches(1.4),
             fill=CHARCOAL, border=None, radius=0.10, shadow=False)
    add_text(slide,
             "« Action Layer » — Klaris qualifie 24/7 en francais quebecois, "
             "execute les relances J+2/J+5, et delegue uniquement la conclusion au courtier.",
             Inches(1.0), Inches(2.30), Inches(SLIDE_W_IN - 2.0), Inches(0.9),
             font=BODY_FONT, size=15, italic=True, color=CREAM, line_spacing=1.35)

    # 3 pillars
    pillars = [
        ("1.", "Toujours dispo",
         "Reponse SMS en < 60 sec, 24/7. Soir, week-end, ferie."),
        ("2.", "Bilingue natif",
         "FR-CA + EN. Detection auto. 'Mise de fonds', 'plex', 'Centris'."),
        ("3.", "Conforme by design",
         "OACIQ + Loi 25 + CASL. Hebergement ca-central-1. Audit logs."),
    ]
    pcard_w = Inches(3.95)
    pcard_h = Inches(2.6)
    gap = Inches(0.18)
    start_x = Inches(0.7)
    y = Inches(3.8)
    for i, (num, head, body) in enumerate(pillars):
        x = start_x + i * (pcard_w + gap)
        add_card(slide, x, y, pcard_w, pcard_h, fill=WHITE,
                 border=BORDER, radius=0.10)
        add_text(slide, num, x + Inches(0.3), y + Inches(0.25),
                 Inches(1), Inches(0.6),
                 font=MONO_FONT, size=26, bold=True, color=TERRACOTTA)
        add_text(slide, head, x + Inches(0.3), y + Inches(0.95),
                 pcard_w - Inches(0.6), Inches(0.5),
                 font=BODY_FONT, size=18, bold=True, color=CHARCOAL,
                 line_spacing=1.15)
        add_text(slide, body, x + Inches(0.3), y + Inches(1.45),
                 pcard_w - Inches(0.6), Inches(1.0),
                 font=BODY_FONT, size=12, color=CHARCOAL_60, line_spacing=1.35)

    # Cost line
    add_text(slide,
             "Cout marginal : ~ 0,05 $ par conversation complete. "
             "Infra : 35-50 $/mois par tenant.",
             Inches(0.7), Inches(6.65), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
             font=MONO_FONT, size=10, color=CHARCOAL_60)

    add_footer(slide, 3, TOTAL_SLIDES)
    set_notes(slide,
              "Killer message v3 etude. 3 piliers = disponibilite + langue + "
              "conformite. Action Layer = positionnement defendable vs CRM "
              "passifs (FUB, kvCORE).")


def slide_04_marche(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Marche adressable",
                       "240 M$ TAM Canada. 17,85 M$ SAM Quebec. SOM cible 600-900 K$.")

    # 3 funnel cards
    funnel = [
        ("TAM", "240 M$", "Canada",
         "160 000 courtiers REALTORS\nARPU mixte 1 500 $/an", CHARCOAL),
        ("SAM", "17,85 M$", "Quebec",
         "~11 900 courtiers OACIQ\nSolos + petites equipes", TERRACOTTA),
        ("SOM", "600-900 K$", "Cible 24 mois",
         "400-600 sieges payants\nCycle B2B 9-18 mois", SUCCESS),
    ]
    fw = Inches(3.95)
    fh = Inches(3.6)
    gap = Inches(0.18)
    start_x = Inches(0.7)
    y = Inches(2.35)
    for i, (lbl, val, geo, sub, col) in enumerate(funnel):
        x = start_x + i * (fw + gap)
        add_card(slide, x, y, fw, fh, fill=WHITE, border=BORDER, radius=0.10)
        # Top stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        x, y, fw, Inches(0.5))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = col
        stripe.line.fill.background()
        add_text(slide, lbl, x + Inches(0.3), y + Inches(0.10),
                 fw - Inches(0.6), Inches(0.35),
                 font=MONO_FONT, size=12, bold=True, color=WHITE)
        add_text(slide, val, x + Inches(0.3), y + Inches(0.8),
                 fw - Inches(0.6), Inches(0.9),
                 font=BODY_FONT, size=34, bold=True, color=CHARCOAL,
                 line_spacing=1.0)
        add_text(slide, geo, x + Inches(0.3), y + Inches(1.85),
                 fw - Inches(0.6), Inches(0.4),
                 font=MONO_FONT, size=11, bold=True, color=col)
        add_text(slide, sub, x + Inches(0.3), y + Inches(2.35),
                 fw - Inches(0.6), Inches(1.2),
                 font=BODY_FONT, size=12, color=CHARCOAL_60, line_spacing=1.45)

    # Source line
    add_text(slide,
             "Sources : CREA Quarterly Forecasts 2024-2025 · OACIQ Statistiques · "
             "APCIQ Barometre · RECO Annual Report 2024 · BCFSA Data Call 2024.",
             Inches(0.7), Inches(6.55), Inches(SLIDE_W_IN - 1.4), Inches(0.5),
             font=MONO_FONT, size=9, italic=True, color=CHARCOAL_60,
             line_spacing=1.35)

    add_footer(slide, 4, TOTAL_SLIDES)
    set_notes(slide,
              "Etude v3 Section 1. SOM 5% du SAM en M24 = penetration "
              "realiste tenant compte cycle B2B 9-18 mois. ARPU mix "
              "70% Solo (1 200 $/an) + 30% Agence (2 400 $/siege/an) = 1 500 $.")


def slide_05_concurrence(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Concurrence",
                       "Moat triangule : conformite + langue + tarif.")

    # Comparison table
    headers = ["Acteur", "Tarif (CAD/mois)", "FR-CA natif", "Loi 25 / CASL", "Position"]
    rows = [
        (["Klaris", "99-199 $", "Oui (Claude natif)", "Oui (ca-central-1)",
          "Action Layer Solo+Agence"], None, TERRACOTTA_SOFT),
        (["Roof.ai (QC)", "Sur devis Entreprise", "Oui", "Oui",
          "B2B Top-down brokerages"], None, None),
        (["Centiva (RE/MAX QC)", "Inclus banniere", "Oui", "Oui",
          "CRM exclusif RE/MAX"], None, None),
        (["Conversica (US)", "~ 4 100 $", "Traduction", "Partielle",
          "Enterprise massif"], None, None),
        (["Structurely (US)", "~ 685 $", "Traduction", "Non garantie",
          "Mid-market US"], None, None),
        (["Follow Up Boss (US)", "80-1 150 $", "Non", "Partielle",
          "CRM passif standard"], None, None),
        (["Adjointe humaine offshore", "2 500-3 000 $", "Limite", "Risque eleve",
          "Substitution couteuse"], None, None),
    ]
    col_widths = [Inches(2.4), Inches(2.2), Inches(2.0), Inches(2.2), Inches(3.1)]
    row_h = Inches(0.42)
    x_start = Inches(0.7)
    y_start = Inches(2.4)
    add_table_row(slide, headers, x_start, y_start, col_widths,
                  row_h, header=True)
    for i, (cells, _, bg) in enumerate(rows):
        add_table_row(slide, cells, x_start,
                      y_start + (i + 1) * row_h,
                      col_widths, row_h, bg=bg)

    # Moat callout
    add_text(slide,
             "Moat = (1) conformite hyper-locale + (2) francais quebecois natif + "
             "(3) tarif destructeur 99 $ vs 500-3 000 $.",
             Inches(0.7), Inches(6.45), Inches(SLIDE_W_IN - 1.4), Inches(0.6),
             font=BODY_FONT, size=12, bold=True, color=TERRACOTTA_STRONG,
             line_spacing=1.3)

    add_footer(slide, 5, TOTAL_SLIDES)
    set_notes(slide,
              "Etude v3 Section 2. Eviter affrontement direct Roof.ai (cycle vente "
              "12 mois top-down). Approche bottom-up via directeurs succursales "
              "independantes. Elise (OACIQ) = menace normative pas concurrentielle.")


def slide_06_produit(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Produit",
                       "App iOS native + miroir web + edge functions Supabase.")

    # Left: feature list
    features = [
        ("SMS qualif", "Twilio + Claude Sonnet · 60 sec · multitour avec memoire"),
        ("Scoring lead", "Grille 6 criteres · froid / tiede / chaud / brulant"),
        ("Relances", "J+2 / J+5 / J+10 · templates approuves courtier"),
        ("Briefing 7h30", "Edge Function cron · email + KPI iOS"),
        ("Voicemail intake", "Whisper FR + Claude · upsert prospect auto"),
        ("Centris sync", "Cron 15 min · listings MLS scopes par agence"),
        ("Watch + Widget", "iOS Live Activities + WatchOS companion"),
        ("Conformite", "Audit logs immuables · export Loi 25 · opt-out CASL"),
    ]
    x = Inches(0.7)
    y0 = Inches(2.4)
    for i, (head, sub) in enumerate(features):
        y = y0 + i * Inches(0.50)
        add_circle(slide, x + Inches(0.10), y + Inches(0.18),
                   Inches(0.16), fill=TERRACOTTA)
        add_text(slide, head, x + Inches(0.35), y, Inches(2.4), Inches(0.4),
                 font=BODY_FONT, size=13, bold=True, color=CHARCOAL,
                 line_spacing=1.2)
        add_text(slide, sub, x + Inches(2.85), y, Inches(4.0), Inches(0.4),
                 font=BODY_FONT, size=11, color=CHARCOAL_60, line_spacing=1.25)

    # Right: phone mockup
    px = Inches(8.4)
    py = Inches(2.0)
    pw = Inches(3.0)
    ph = Inches(5.0)
    phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, py, pw, ph)
    phone.adjustments[0] = 0.10
    phone.fill.solid()
    phone.fill.fore_color.rgb = CHARCOAL
    phone.line.color.rgb = CHARCOAL
    phone.line.width = Pt(2)

    # Screen
    sx = px + Inches(0.18)
    sy = py + Inches(0.4)
    sw = pw - Inches(0.36)
    sh = ph - Inches(0.8)
    screen = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    sx, sy, sw, sh)
    screen.adjustments[0] = 0.05
    screen.fill.solid()
    screen.fill.fore_color.rgb = CREAM
    screen.line.fill.background()

    # Screen content: header + 3 prospect cards
    add_text(slide, "Aujourd'hui · 3 chauds",
             sx + Inches(0.2), sy + Inches(0.25), sw - Inches(0.4), Inches(0.3),
             font=BODY_FONT, size=10, bold=True, color=CHARCOAL)
    add_text(slide, "Klaris · 8 leads traites",
             sx + Inches(0.2), sy + Inches(0.5), sw - Inches(0.4), Inches(0.3),
             font=MONO_FONT, size=8, color=CHARCOAL_60)

    prospects = [
        ("Marie L.", "Plateau · 645 K$", 5, DESTRUCTIVE),
        ("Karim S.", "Laval · 680 K$", 4, TERRACOTTA),
        ("Ann D.", "Verdun · 595 K$", 3, WARN_AMBER),
    ]
    for i, (name, info, heat, col) in enumerate(prospects):
        cy = sy + Inches(0.95) + i * Inches(0.85)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      sx + Inches(0.15), cy,
                                      sw - Inches(0.3), Inches(0.7))
        card.adjustments[0] = 0.12
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER
        card.line.width = Pt(0.5)
        add_circle(slide, sx + Inches(0.35), cy + Inches(0.35),
                   Inches(0.30), fill=col)
        add_text(slide, name,
                 sx + Inches(0.6), cy + Inches(0.08),
                 sw - Inches(0.8), Inches(0.28),
                 font=BODY_FONT, size=10, bold=True, color=CHARCOAL,
                 line_spacing=1.1)
        add_text(slide, info,
                 sx + Inches(0.6), cy + Inches(0.35),
                 sw - Inches(0.8), Inches(0.28),
                 font=MONO_FONT, size=8, color=CHARCOAL_60, line_spacing=1.1)

    add_footer(slide, 6, TOTAL_SLIDES)
    set_notes(slide,
              "Produit live: Sprint 1-6 livres. Stack n8n + Claude Sonnet + "
              "Supabase + Twilio. App Flutter/Cupertino iOS. Cout infra "
              "35-50 $/tenant/mois.")


def slide_07_traction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Traction",
                       "Pilote actif · 2 beta-testeurs · feedback Q1-Q2 2026.")

    # 2 beta cards
    betas = [
        ("Joanel", "Courtier independant",
         "Grand Montreal · bilingue · ~ 2-4 M$ GCI",
         "Validation produit-marche Solo. Citation hero deck."),
        ("Maxime Belma", "RE/MAX Anjou — directeur",
         "Banniere RE/MAX · Persona 4 · 50+ courtiers",
         "Validation Persona Agence + cohabitation Centiva (Action Layer)."),
    ]
    cw = Inches(5.95)
    ch = Inches(3.0)
    gap = Inches(0.30)
    start_x = Inches(0.7)
    y = Inches(2.4)
    for i, (name, role, ctx, learning) in enumerate(betas):
        x = start_x + i * (cw + gap)
        add_card(slide, x, y, cw, ch, fill=WHITE, border=BORDER, radius=0.10)
        add_text(slide, name, x + Inches(0.3), y + Inches(0.25),
                 cw - Inches(0.6), Inches(0.5),
                 font=BODY_FONT, size=22, bold=True, color=CHARCOAL,
                 line_spacing=1.1)
        add_text(slide, role, x + Inches(0.3), y + Inches(0.85),
                 cw - Inches(0.6), Inches(0.4),
                 font=MONO_FONT, size=11, bold=True, color=TERRACOTTA)
        add_text(slide, ctx, x + Inches(0.3), y + Inches(1.30),
                 cw - Inches(0.6), Inches(0.5),
                 font=BODY_FONT, size=12, color=CHARCOAL_60, line_spacing=1.3)
        add_text(slide, "Apprentissage cle :",
                 x + Inches(0.3), y + Inches(1.95),
                 cw - Inches(0.6), Inches(0.3),
                 font=MONO_FONT, size=9, bold=True, color=CHARCOAL_60)
        add_text(slide, learning, x + Inches(0.3), y + Inches(2.25),
                 cw - Inches(0.6), Inches(0.7),
                 font=BODY_FONT, size=13, italic=True, color=CHARCOAL,
                 line_spacing=1.35)

    # KPI strip
    kpi_y = Inches(5.6)
    kw = Inches(3.85)
    kh = Inches(1.3)
    kpis = [
        ("Sprints livres", "1-6", "iOS + Web + Edge Functions"),
        ("Statut pilote", "Live", "Q2 2026 · feedback en cours"),
        ("Conversion cible M6", "50 payants", "Jalon product-market fit"),
    ]
    for i, (lbl, val, sub) in enumerate(kpis):
        x = Inches(0.7) + i * (kw + Inches(0.20))
        add_kpi_card(slide, x, kpi_y, kw, kh, lbl, val, sub,
                     value_color=TERRACOTTA)

    add_footer(slide, 7, TOTAL_SLIDES)
    set_notes(slide,
              "Traction reelle: 2 beta-testeurs nommes (Joanel + Belma). "
              "Pas de chiffres conversations/leads encore consolides — placeholder "
              "honnete 'feedback en cours'. A mettre a jour avant pitch live.")


def slide_08_business_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Modele d'affaires",
                       "Trois forfaits. ARPU mix 1 500 $/an. SaaS pur, pas success-fee.")

    # 3 pricing tiers
    tiers = [
        ("Solo Starter", "99 $", "/mois",
         ["50 conversations IA/mois", "Numero local Twilio",
          "Briefing 7h30 quotidien", "CRM leger + scoring"],
         WHITE, CHARCOAL),
        ("Solo Pro", "149 $", "/mois",
         ["Conversations illimitees", "Outbound proactif J+5/J+15",
          "Whisper notes vocales illimite", "API Follow Up Boss + Centiva"],
         CHARCOAL, CREAM),
        ("Agence", "199 $", "/siege/mois",
         ["Min 5 sieges · SLA 99,9 %", "Manager Dashboard + RLS",
          "Kit conformite Loi 25 (DPO)", "Audit logs + rapport PDF mensuel"],
         TERRACOTTA, WHITE),
    ]
    tw = Inches(3.95)
    th = Inches(3.6)
    gap = Inches(0.18)
    start_x = Inches(0.7)
    y = Inches(2.35)
    for i, (name, price, per, feats, fill, txt_color) in enumerate(tiers):
        x = start_x + i * (tw + gap)
        add_card(slide, x, y, tw, th, fill=fill, border=BORDER, radius=0.12)
        add_text(slide, name, x + Inches(0.3), y + Inches(0.3),
                 tw - Inches(0.6), Inches(0.5),
                 font=BODY_FONT, size=18, bold=True, color=txt_color)
        # Price
        add_text(slide, price, x + Inches(0.3), y + Inches(0.85),
                 Inches(2.0), Inches(0.9),
                 font=BODY_FONT, size=42, bold=True,
                 color=TERRACOTTA if fill != TERRACOTTA else WHITE,
                 line_spacing=1.0)
        add_text(slide, per, x + Inches(2.1), y + Inches(1.40),
                 Inches(1.8), Inches(0.3),
                 font=MONO_FONT, size=11, color=txt_color)
        # Features
        for j, feat in enumerate(feats):
            fy = y + Inches(1.95) + j * Inches(0.35)
            add_circle(slide, x + Inches(0.40), fy + Inches(0.13),
                       Inches(0.12),
                       fill=TERRACOTTA if fill != TERRACOTTA else WHITE)
            add_text(slide, feat, x + Inches(0.6), fy,
                     tw - Inches(0.8), Inches(0.32),
                     font=BODY_FONT, size=11, color=txt_color, line_spacing=1.25)

    # ARPU strip
    add_text(slide,
             "ARPU mix 70 % Solo + 30 % Agence = 1 500 $/an. "
             "Reference comparative : adjointe humaine 2 500-3 000 $/mois.",
             Inches(0.7), Inches(6.55), Inches(SLIDE_W_IN - 1.4), Inches(0.5),
             font=MONO_FONT, size=10, color=CHARCOAL_60, line_spacing=1.4)

    add_footer(slide, 8, TOTAL_SLIDES)
    set_notes(slide,
              "Etude v3 Section 3. Pas de success-fee (illegal OACIQ sans permis "
              "courtage). Agence engagement 5 sieges min = upsell mechanic NRR.")


def slide_09_unit_economics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Unit economics",
                       "CAC 1 200 $ · LTV 3 000 $ · Payback 9,6 mois · NRR > 100 %.")

    kpis = [
        ("CAC numerique", "1 200 $", "Meta + LinkedIn · cible courtiers OACIQ",
         TERRACOTTA),
        ("LTV (24 mois)", "3 000 $", "ARPU 1 500 $/an × retention 2 ans",
         SUCCESS),
        ("Ratio LTV / CAC", "2,5 : 1", "Acceptable Bootstrap · standard SaaS 3:1",
         WARN_AMBER),
        ("Payback period", "9,6 mois", "MRR moyen 125 $ · cash drawdown 12 mois",
         INFO),
    ]
    kw = Inches(2.95)
    kh = Inches(2.0)
    gap = Inches(0.15)
    start_x = Inches(0.7)
    y = Inches(2.4)
    for i, (lbl, val, sub, col) in enumerate(kpis):
        x = start_x + i * (kw + gap)
        add_kpi_card(slide, x, y, kw, kh, lbl, val, sub, value_color=col)

    # Bottom strip — assumptions
    y2 = Inches(4.85)
    add_card(slide, Inches(0.7), y2,
             Inches(SLIDE_W_IN - 1.4), Inches(2.1),
             fill=WHITE, border=BORDER, radius=0.10)
    add_text(slide, "HYPOTHESES & LEVIERS",
             Inches(0.95), y2 + Inches(0.20), Inches(8), Inches(0.3),
             font=MONO_FONT, size=10, bold=True, color=CHARCOAL_60)

    levers = [
        (SUCCESS, "Churn cible", "< 3 %/mois (norme PME SaaS B2B)"),
        (SUCCESS, "NRR cible 24 mois", "> 100 % via upsell Solo Pro → Agence"),
        (WARN_AMBER, "Referral program", "100 $/100 $ · baisse CAC blended"),
        (INFO, "Cout marginal/conv.", "~ 0,05 $ · infra 35-50 $/tenant/mois"),
    ]
    for i, (col, lbl, val) in enumerate(levers):
        row_y = y2 + Inches(0.6) + i * Inches(0.34)
        add_bullet_line(slide, col, lbl, val,
                        Inches(0.95), row_y, Inches(SLIDE_W_IN - 1.9))

    add_footer(slide, 9, TOTAL_SLIDES)
    set_notes(slide,
              "Etude v3 Section 4. CAC 1 200 $ realiste audience OACIQ. "
              "Payback 9,6 mois limite haute. Action critique: Tech E&O quote "
              "(Northbridge / Intact Tech) avant Go-Live = 5-10 K$/an.")


def slide_10_gtm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Go-To-Market",
                       "3 phases sequencees. 5 partenariats fondateurs annee 1.")

    # 3 phases timeline
    phases = [
        ("Phase 1", "M1-M12", "Focus Quebec",
         "RMR Montreal · Persona Solo bilingue · Bottom-up",
         TERRACOTTA, "GO"),
        ("Phase 2", "M12-M24", "Ontario + C.-B.",
         "Conditionnel cash-flow positif QC · TRESA · MLS local",
         WARN_AMBER, "GO conditionnel"),
        ("Phase 3", "M24+", "US frontaliers",
         "Floride FR/ES · Louisiane · Nouveau-Brunswick · Mexique",
         CHARCOAL_60, "NO-GO actuel"),
    ]
    pw = Inches(3.95)
    ph = Inches(2.5)
    gap = Inches(0.18)
    start_x = Inches(0.7)
    y = Inches(2.4)
    for i, (lbl, dur, geo, sub, col, verdict) in enumerate(phases):
        x = start_x + i * (pw + gap)
        add_card(slide, x, y, pw, ph, fill=WHITE, border=BORDER, radius=0.10)
        # Verdict pill
        vpill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x + Inches(0.3), y + Inches(0.3),
                                       Inches(1.8), Inches(0.32))
        vpill.adjustments[0] = 0.5
        vpill.fill.solid()
        vpill.fill.fore_color.rgb = col
        vpill.line.fill.background()
        add_text(slide, verdict, x + Inches(0.3), y + Inches(0.32),
                 Inches(1.8), Inches(0.3),
                 font=MONO_FONT, size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(slide, lbl, x + Inches(0.3), y + Inches(0.75),
                 pw - Inches(0.6), Inches(0.4),
                 font=BODY_FONT, size=18, bold=True, color=CHARCOAL)
        add_text(slide, dur, x + Inches(0.3), y + Inches(1.10),
                 pw - Inches(0.6), Inches(0.3),
                 font=MONO_FONT, size=10, color=CHARCOAL_60)
        add_text(slide, geo, x + Inches(0.3), y + Inches(1.40),
                 pw - Inches(0.6), Inches(0.4),
                 font=BODY_FONT, size=14, bold=True, color=col)
        add_text(slide, sub, x + Inches(0.3), y + Inches(1.80),
                 pw - Inches(0.6), Inches(0.8),
                 font=BODY_FONT, size=11, color=CHARCOAL_60, line_spacing=1.35)

    # 5 partnerships strip
    pn_y = Inches(5.25)
    add_card(slide, Inches(0.7), pn_y,
             Inches(SLIDE_W_IN - 1.4), Inches(1.75),
             fill=TERRACOTTA_SOFT, border=None, radius=0.10, shadow=False)
    add_text(slide, "5 PARTENARIATS FONDATEURS — ANNEE 1",
             Inches(0.95), pn_y + Inches(0.18), Inches(8), Inches(0.3),
             font=MONO_FONT, size=10, bold=True, color=TERRACOTTA_STRONG)
    partners = [
        "1. Centris Inc. — entente syndication API MLS",
        "2. College de l'immobilier QC — cursus UFC technologies",
        "3. APCIQ — commandite RDV OACIQ + evenements statistiques",
        "4. Succursales franchisees innovantes — RE/MAX Anjou (Belma) pilote",
        "5. Cabinet juridique Loi 25 — livre blanc co-redige (lead magnet)",
    ]
    for i, p in enumerate(partners):
        col = 0 if i < 3 else 1
        row = i if i < 3 else i - 3
        px = Inches(0.95) + col * Inches(6.0)
        py = pn_y + Inches(0.55) + row * Inches(0.32)
        add_text(slide, p, px, py, Inches(6), Inches(0.3),
                 font=BODY_FONT, size=11, color=CHARCOAL_60, line_spacing=1.2)

    add_footer(slide, 10, TOTAL_SLIDES)
    set_notes(slide,
              "Etude v3 Section 7. Strategie bottom-up sur Persona 4 (directeurs "
              "franchisees) plutot que top-down sieges sociaux. Cycle vente "
              "agence 9-18 mois.")


def slide_11_risques(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Risques majeurs & mitigation",
                       "3 risques operationnels critiques. Defenses cablees.")

    headers = ["Risque", "Impact", "Mitigation"]
    rows = [
        (["1. Hallucinations LLM (fausse info SMS)",
          "Critique — responsabilite courtier OACIQ + FARCIQ",
          "RAG strict · Temperature = 0 · ToS Limitation Liability · Tech E&O Klaris 5-10 K$/an"],
         DESTRUCTIVE),
        (["2. Blocage SMS Twilio (CRTC / A2P)",
          "Eleve — coupe valeur produit (suivis non livres)",
          "Toll-Free Verification CA · Sender ID Bell/Rogers/Telus · Double Opt-in CASL"],
         WARN_AMBER),
        (["3. Revocation acces Centris API",
          "Majeur — reco listings degradees",
          "Entente syndication Centris · Fallback DDF CREA broker-level · pont technique"],
         INFO),
    ]
    col_widths = [Inches(4.0), Inches(3.8), Inches(4.1)]
    row_h = Inches(0.95)
    x_start = Inches(0.7)
    y_start = Inches(2.4)
    add_table_row(slide, headers, x_start, y_start, col_widths,
                  Inches(0.5), header=True)
    for i, (cells, col_dot) in enumerate(rows):
        add_table_row(slide, cells, x_start,
                      y_start + Inches(0.5) + i * row_h,
                      col_widths, row_h)
        # Left color marker
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        x_start, y_start + Inches(0.5) + i * row_h,
                                        Inches(0.10), row_h)
        marker.fill.solid()
        marker.fill.fore_color.rgb = col_dot
        marker.line.fill.background()

    # Bottom callout
    add_text(slide,
             "Le cadre OACIQ + Loi 25 + CASL = fosse defensif. "
             "Concurrent US doit refondre architecture pour entrer au QC.",
             Inches(0.7), Inches(6.4), Inches(SLIDE_W_IN - 1.4), Inches(0.6),
             font=BODY_FONT, size=12, bold=True, italic=True, color=TERRACOTTA_STRONG,
             line_spacing=1.3)

    add_footer(slide, 11, TOTAL_SLIDES)
    set_notes(slide,
              "Etude v3 Section 5. Tech E&O au Quebec PME tech = 3-12 K$/an. "
              "Budget 5-10 K$. Action urgente Q1: quote ferme Northbridge / "
              "Intact Tech / Assur360.")


def slide_12_roadmap_kpi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_section_header(slide, "Roadmap & KPI North-Star",
                       "Visites qualifiees autonomes / semaine. 3 jalons sur 24 mois.")

    # North-star banner
    add_card(slide, Inches(0.7), Inches(2.0),
             Inches(SLIDE_W_IN - 1.4), Inches(1.0),
             fill=CHARCOAL, border=None, radius=0.10, shadow=False)
    add_text(slide, "NORTH-STAR METRIC",
             Inches(0.95), Inches(2.15), Inches(8), Inches(0.3),
             font=MONO_FONT, size=10, bold=True, color=CREAM)
    add_text(slide,
             "Nombre de rencontres physiques qualifiees confirmees "
             "de maniere autonome par Klaris · base hebdomadaire.",
             Inches(0.95), Inches(2.45), Inches(SLIDE_W_IN - 2.0), Inches(0.5),
             font=BODY_FONT, size=14, color=CREAM, line_spacing=1.3)

    # 3 milestones
    miles = [
        ("M6", "50 clients", "60 K$ ARR",
         "Product-market fit · retention nette > 90 %", TERRACOTTA),
        ("M12", "100 clients", "150 K$ ARR",
         "Equipe agence 1ere · 10+ courtiers integres", WARN_AMBER),
        ("M24", "400 clients", "600 K$ ARR",
         "Cash-flow positif · trigger Phase 2 Ontario", SUCCESS),
    ]
    mw = Inches(3.95)
    mh = Inches(2.8)
    gap = Inches(0.18)
    start_x = Inches(0.7)
    y = Inches(3.4)
    for i, (jalon, clients, arr, sub, col) in enumerate(miles):
        x = start_x + i * (mw + gap)
        add_card(slide, x, y, mw, mh, fill=WHITE, border=BORDER, radius=0.10)
        add_text(slide, jalon, x + Inches(0.3), y + Inches(0.25),
                 mw - Inches(0.6), Inches(0.45),
                 font=MONO_FONT, size=14, bold=True, color=col)
        add_text(slide, clients, x + Inches(0.3), y + Inches(0.75),
                 mw - Inches(0.6), Inches(0.7),
                 font=BODY_FONT, size=30, bold=True, color=CHARCOAL,
                 line_spacing=1.0)
        add_text(slide, arr, x + Inches(0.3), y + Inches(1.55),
                 mw - Inches(0.6), Inches(0.5),
                 font=BODY_FONT, size=18, bold=True, color=TERRACOTTA,
                 line_spacing=1.1)
        add_text(slide, sub, x + Inches(0.3), y + Inches(2.05),
                 mw - Inches(0.6), Inches(0.7),
                 font=BODY_FONT, size=11, color=CHARCOAL_60, line_spacing=1.35)

    # KPI to watch
    add_text(slide,
             "Metriques surveillees : retention nette mensuelle > 90 % · NRR > 100 % · "
             "churn logo < 3 %/mois · payback < 12 mois.",
             Inches(0.7), Inches(6.55), Inches(SLIDE_W_IN - 1.4), Inches(0.5),
             font=MONO_FONT, size=10, italic=True, color=CHARCOAL_60,
             line_spacing=1.35)

    add_footer(slide, 12, TOTAL_SLIDES)
    set_notes(slide,
              "Etude v3 Section 7.4. North-Star = anti-vanity (pas downloads / "
              "visites). M12 100 clients = jalon PMF. M24 400 = trigger Phase 2 "
              "Ontario (cash-flow positif).")


def slide_13_demande_equipe(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CREAM)
    add_grid_pattern(slide, Inches(0), Inches(0),
                     Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
                     step_in=0.45, alpha=0.30)

    add_section_header(slide, "Equipe & demande",
                       "Bootstrap. Ce qu'on cherche : intros, pas chequier.")

    # Team — 4 cofounders
    members = [
        ("Dennis", "Co-founder · Tech lead",
         "n8n · Claude API · Supabase · Edge Functions"),
        ("Eliot", "Co-founder · Strategie",
         "Marche QC · partenariats · ops"),
        ("Walkens", "Co-founder",
         "Produit + finance · ops courtage"),
        ("Seydou", "Co-founder",
         "Reseau brokerages QC · go-to-market"),
    ]
    mw = Inches(2.95)
    mh = Inches(1.5)
    gap = Inches(0.15)
    start_x = Inches(0.7)
    y = Inches(2.3)
    for i, (name, role, sub) in enumerate(members):
        x = start_x + i * (mw + gap)
        add_card(slide, x, y, mw, mh, fill=WHITE, border=BORDER, radius=0.10)
        # Avatar circle
        add_circle(slide, x + Inches(0.45), y + Inches(0.45),
                   Inches(0.55), fill=TERRACOTTA_SOFT)
        add_text(slide, name[0], x + Inches(0.20), y + Inches(0.22),
                 Inches(0.5), Inches(0.5),
                 font=BODY_FONT, size=20, bold=True, color=TERRACOTTA_STRONG,
                 align=PP_ALIGN.CENTER)
        add_text(slide, name, x + Inches(0.85), y + Inches(0.20),
                 mw - Inches(1.0), Inches(0.4),
                 font=BODY_FONT, size=14, bold=True, color=CHARCOAL)
        add_text(slide, role, x + Inches(0.85), y + Inches(0.55),
                 mw - Inches(1.0), Inches(0.35),
                 font=MONO_FONT, size=9, color=TERRACOTTA)
        add_text(slide, sub, x + Inches(0.30), y + Inches(0.95),
                 mw - Inches(0.6), Inches(0.5),
                 font=BODY_FONT, size=10, color=CHARCOAL_60, line_spacing=1.3)

    # Critical open role
    add_card(slide, Inches(0.7), Inches(4.05),
             Inches(SLIDE_W_IN - 1.4), Inches(1.0),
             fill=WARN_AMBER, border=None, radius=0.10, shadow=False)
    add_text(slide, "DECISION Q1 OUVERTE",
             Inches(0.95), Inches(4.20), Inches(8), Inches(0.3),
             font=MONO_FONT, size=10, bold=True, color=CHARCOAL)
    add_text(slide,
             "Assignation VP Ventes (1 cofondateur a 100 %) — cycle B2B agence "
             "9-18 mois · pipeline Persona 4. Sans ce role, jalon M12 = 100 clients ne livre pas.",
             Inches(0.95), Inches(4.50), Inches(SLIDE_W_IN - 2.0), Inches(0.55),
             font=BODY_FONT, size=12, color=CHARCOAL, line_spacing=1.3)

    # Asks
    add_card(slide, Inches(0.7), Inches(5.30),
             Inches(SLIDE_W_IN - 1.4), Inches(1.65),
             fill=CHARCOAL, border=None, radius=0.10, shadow=False)
    add_text(slide, "CE QU'ON CHERCHE",
             Inches(0.95), Inches(5.45), Inches(8), Inches(0.3),
             font=MONO_FONT, size=10, bold=True, color=TERRACOTTA_SOFT)
    asks = [
        "1. Intro Centris Inc. — direction innovation / partenariats API.",
        "2. Intro RE/MAX / Royal LePage / Sutton QC — direction franchise tech.",
        "3. Intro OACIQ direction strategie — alignement Elise + standards IA.",
        "4. Avocat senior droit techno + Loi 25 — co-redaction livre blanc.",
    ]
    for i, ask in enumerate(asks):
        col = 0 if i < 2 else 1
        row = i if i < 2 else i - 2
        ax = Inches(0.95) + col * Inches(6.0)
        ay = Inches(5.80) + row * Inches(0.35)
        add_text(slide, ask, ax, ay, Inches(5.8), Inches(0.32),
                 font=BODY_FONT, size=11, color=CREAM, line_spacing=1.3)

    # Contact line
    add_text(slide,
             "alanmanou.consulting@gmail.com  ·  klaris.app  ·  Montreal, QC",
             Inches(0.7), Inches(7.1), Inches(SLIDE_W_IN - 1.4), Inches(0.3),
             font=MONO_FONT, size=10, bold=True, color=CHARCOAL_60,
             align=PP_ALIGN.CENTER)

    add_footer(slide, 13, TOTAL_SLIDES)
    set_notes(slide,
              "Close: bootstrap explicite. Demande = intros, pas equity. "
              "VP Sales role ouvert Q1 = honnetete sur execution. "
              "Annexes data room: stress test churn, P&L Y1-Y3, cap table, "
              "ToS draft, etude v3 complete.")


# ----------------------------------------------------------------------------
# Build
# ----------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slide_01_cover(prs)
    slide_02_probleme(prs)
    slide_03_solution(prs)
    slide_04_marche(prs)
    slide_05_concurrence(prs)
    slide_06_produit(prs)
    slide_07_traction(prs)
    slide_08_business_model(prs)
    slide_09_unit_economics(prs)
    slide_10_gtm(prs)
    slide_11_risques(prs)
    slide_12_roadmap_kpi(prs)
    slide_13_demande_equipe(prs)

    from klaris_annexes import slide_annexe_acronymes, slide_annexe_metriques
    slide_annexe_acronymes(prs, 14, TOTAL_SLIDES)
    slide_annexe_metriques(prs, 15, TOTAL_SLIDES)

    out = Path(__file__).parent / "Klaris_Investisseur_FR.pptx"
    prs.save(out)
    print(f"OK -> {out}")


if __name__ == "__main__":
    build()
